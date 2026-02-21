"""
ModelIt Lesson Materials Generator v4
Creates PowerPoint, Student Activity Pack (Word), and Teacher's Guide (Word)
with AI-generated images via NanoBanana (OpenRouter API).

v4 Improvements:
- Proper image aspect ratio preservation (no warping)
- Enhanced intro slides with more content
- Comprehensive Teacher's Guide with:
  - CAST testing information
  - Background content knowledge
  - Slide-by-slide facilitation guide
  - Expected student responses
- Restored Student Activity Pack with open response questions
- Typography and accessibility best practices
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGB, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import requests
import base64
import os
import time
from PIL import Image

# ============================================
# CONFIGURATION
# ============================================
LESSON_ID = "G05-L01"
LESSON_TITLE = "When Trees Become Matches"
LESSON_SUBTITLE = "California's Burning Season and the Earth Systems That Fuel It"
NGSS_STANDARD = "5-ESS2-1"
GRADE_LEVEL = "5th Grade"
AGE_RANGE = "10-11 years old"

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'materials', 'grade-05', LESSON_ID)
IMAGES_DIR = os.path.join(OUTPUT_DIR, 'images')

# API Configuration
OPENROUTER_API_KEY = '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'

# ============================================
# BRAND COLORS
# ============================================
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
BRAND_BLUE = RGBColor(0x1A, 0x47, 0x80)
MID_BLUE = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_BLUE = RGBColor(0x7E, 0xC8, 0xE3)
SKY_BLUE = RGBColor(0x5D, 0xB7, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)

# Word colors
DOCX_NAVY = DocxRGB(0x0D, 0x1B, 0x2A)
DOCX_BRAND_BLUE = DocxRGB(0x1A, 0x47, 0x80)
DOCX_MID_BLUE = DocxRGB(0x2E, 0x86, 0xAB)
DOCX_LIGHT_BLUE = DocxRGB(0x7E, 0xC8, 0xE3)
DOCX_ORANGE = DocxRGB(0xE6, 0x7E, 0x22)

# Cost tracking
total_image_cost = 0.0

def ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(IMAGES_DIR, exist_ok=True)


# ============================================
# IMAGE GENERATION (NanoBanana via OpenRouter)
# ============================================

def generate_image(scene_description, filename, setting="modern classroom"):
    """
    Generate a photorealistic educational image using NanoBanana.
    Uses the approved prompt template for diversity and style.
    """
    global total_image_cost

    prompt = f'''Create a photorealistic, cinematic, ultra-crisp image of {scene_description} in {setting} featuring a diverse, multicultural group with Black people centered (a mix of skin tones and ethnicities represented), age-accurate for {AGE_RANGE} (no one looks like an adult if they're a kid).

Style: modern education / future-ready / "middle school coolness"--confident, aspirational, realistic.

Composition: clean framing, strong focal subject, natural body proportions, believable clothing textures, realistic hair detail (coils, curls, locs, braids), subtle emotion and authentic interaction.

Camera/lighting: soft natural window light + gentle fill, shallow depth of field, 35mm/50mm look, high dynamic range, sharp eyes, realistic skin texture.

Quality: high-resolution, crisp edges, minimal noise, professional editorial photo.

If tech overlays appear: keep them subtle, readable, and physically plausible (not cluttered), with realistic reflections and occlusion.'''

    print(f"  Generating: {filename}...")

    try:
        response = requests.post(
            'https://openrouter.ai/api/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {OPENROUTER_API_KEY}',
                'Content-Type': 'application/json',
            },
            json={
                'model': 'google/gemini-2.5-flash-image',
                'messages': [{'role': 'user', 'content': prompt}],
                'modalities': ['image', 'text']
            },
            timeout=180
        )

        if response.status_code == 200:
            data = response.json()
            choices = data.get('choices', [])

            if choices:
                message = choices[0].get('message', {})
                images = message.get('images', [])

                if images:
                    img = images[0]
                    if isinstance(img, dict):
                        img_url = img.get('image_url', {}).get('url', '')
                        if img_url.startswith('data:image'):
                            base64_data = img_url.split(',')[1]
                            filepath = os.path.join(IMAGES_DIR, filename)
                            with open(filepath, 'wb') as f:
                                f.write(base64.b64decode(base64_data))

                            # Track cost
                            usage = data.get('usage', {})
                            cost = usage.get('cost', 0)
                            total_image_cost += cost
                            print(f"    [OK] Saved: {filename} (${cost:.4f})")
                            return filepath

            print(f"    [!] No image in response for {filename}")
            return None
        else:
            print(f"    [X] Error {response.status_code}: {response.text[:200]}")
            return None

    except Exception as e:
        print(f"    [X] Exception: {e}")
        return None


def generate_all_lesson_images():
    """Generate all images needed for the lesson materials."""
    print("\n" + "=" * 60)
    print("GENERATING LESSON IMAGES")
    print("=" * 60)

    images = {}

    # Image 1: Cover - Students collaborating on wildfire project
    images['cover'] = generate_image(
        "5th grade students collaborating on a science modeling project about California wildfires, gathered around a laptop showing a diagram, one student confidently pointing at the screen while others lean in engaged",
        "cover-students-collaborating.png",
        "modern classroom with natural lighting"
    )
    time.sleep(2)  # Rate limiting

    # Image 2: California landscape - fire season
    images['landscape'] = generate_image(
        "California golden hillside during fire season with dry grass and scattered oak trees, warm sunset lighting, a diverse suburban community visible in the distance, educational documentary style",
        "california-fire-season.png",
        "outdoor California landscape"
    )
    time.sleep(2)

    # Image 3: Students discussing/comparing
    images['discussion'] = generate_image(
        "a small group of 5th grade students having an animated science discussion, some pointing at a shared screen, others taking notes, showing genuine curiosity and engagement, one student explaining to others",
        "students-discussion.png",
        "bright modern classroom or maker space"
    )
    time.sleep(2)

    # Image 4: Student working on model
    images['modeling'] = generate_image(
        "a focused 5th grade student working on a laptop showing a simple diagram with connected nodes and arrows, the student looks confident and engaged, other students visible in background collaborating",
        "student-modeling.png",
        "modern classroom with Chromebooks"
    )
    time.sleep(2)

    # Image 5: STEM Challenge - engineering design
    images['stem'] = generate_image(
        "5th grade students working together on an engineering design challenge, sketching on paper and discussing ideas, looking like young engineers solving a real problem, hands-on collaborative work",
        "stem-challenge.png",
        "classroom maker space with materials"
    )

    print(f"\nTotal images generated: {len([v for v in images.values() if v])}")
    print(f"Total image cost: ${total_image_cost:.4f}")

    return images


# ============================================
# IMAGE UTILITIES - ASPECT RATIO PRESERVATION
# ============================================

def get_image_dimensions(image_path):
    """Get image dimensions using PIL."""
    try:
        with Image.open(image_path) as img:
            return img.width, img.height
    except Exception as e:
        print(f"    [!] Could not read image dimensions: {e}")
        return None, None


def calculate_scaled_dimensions(image_path, max_width, max_height):
    """
    Calculate scaled dimensions that preserve aspect ratio.
    Returns (width, height) in inches that fit within max bounds.
    """
    orig_width, orig_height = get_image_dimensions(image_path)
    if not orig_width or not orig_height:
        return max_width, max_height  # Fallback to max dimensions

    aspect_ratio = orig_width / orig_height

    # Try fitting to width first
    scaled_width = max_width
    scaled_height = max_width / aspect_ratio

    # If height exceeds max, scale to fit height instead
    if scaled_height > max_height:
        scaled_height = max_height
        scaled_width = max_height * aspect_ratio

    return scaled_width, scaled_height


# ============================================
# POWERPOINT FUNCTIONS
# ============================================

def add_diagonal_corners(slide, prs, position='both'):
    """Add the signature blue diagonal corner design"""
    if position in ['top', 'both']:
        shape1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(-0.5), Inches(-0.5), Inches(4), Inches(1.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = NAVY
        shape1.line.fill.background()

        shape2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(2.5), Inches(0), Inches(3), Inches(0.8))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = MID_BLUE
        shape2.line.fill.background()

        shape3 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(3.5), Inches(0.5), Inches(2.5), Inches(0.4))
        shape3.fill.solid()
        shape3.fill.fore_color.rgb = LIGHT_BLUE
        shape3.line.fill.background()

    if position in ['bottom', 'both']:
        shape4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(6.5), Inches(6), Inches(4), Inches(1.5))
        shape4.fill.solid()
        shape4.fill.fore_color.rgb = NAVY
        shape4.line.fill.background()
        shape4.rotation = 180

        shape5 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(4.5), Inches(6.7), Inches(3), Inches(0.8))
        shape5.fill.solid()
        shape5.fill.fore_color.rgb = MID_BLUE
        shape5.line.fill.background()

        shape6 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(4), Inches(6.3), Inches(2.5), Inches(0.4))
        shape6.fill.solid()
        shape6.fill.fore_color.rgb = LIGHT_BLUE
        shape6.line.fill.background()


def add_modelit_logo(slide, x, y, width=3):
    """Add ModelIt! logo box"""
    logo_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(0.9))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = SKY_BLUE
    logo_bg.line.fill.background()

    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x + 0.15), Inches(y + 0.25), Inches(0.4), Inches(0.4))
    icon.fill.solid()
    icon.fill.fore_color.rgb = WHITE
    icon.line.fill.background()

    text_box = slide.shapes.add_textbox(
        Inches(x + 0.6), Inches(y + 0.15), Inches(width - 0.8), Inches(0.6))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ModelIt!"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Comic Sans MS"


def add_image_to_slide(slide, image_path, x, y, max_width, max_height, center=False):
    """
    Add an image to a slide while PRESERVING ASPECT RATIO.
    No more warping/stretching - images maintain their natural proportions.
    If center=True, centers the image within the max bounds.
    """
    if image_path and os.path.exists(image_path):
        # Calculate dimensions that preserve aspect ratio
        scaled_width, scaled_height = calculate_scaled_dimensions(
            image_path, max_width, max_height
        )

        # Optionally center within the max bounds
        if center:
            x_offset = (max_width - scaled_width) / 2
            y_offset = (max_height - scaled_height) / 2
            x = x + x_offset
            y = y + y_offset

        slide.shapes.add_picture(
            image_path,
            Inches(x),
            Inches(y),
            Inches(scaled_width),
            Inches(scaled_height)
        )
        return True
    else:
        # Fallback placeholder
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(max_width), Inches(max_height))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFA)
        box.line.color.rgb = MID_BLUE
        box.line.width = Pt(2)
        return False


def add_slide_number(slide, number, total=8):
    """Add slide number indicator"""
    text_box = slide.shapes.add_textbox(
        Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.RIGHT


def create_powerpoint(images):
    """Create the student presentation PowerPoint with proper aspect ratios."""
    output_path = os.path.join(OUTPUT_DIR, f'{LESSON_ID}-Student-Presentation.pptx')
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: COVER ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs)
    add_modelit_logo(slide, 3.5, 1.3)

    # Add cover image with proper aspect ratio
    if images.get('cover'):
        add_image_to_slide(slide, images['cover'], 0.3, 4.5, 3.2, 2.5, center=True)

    header = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Student Lesson"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    title = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = LESSON_TITLE
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = LESSON_SUBTITLE
    sub.font.size = Pt(15)
    sub.font.italic = True
    sub.font.color.rgb = DARK_TEXT
    sub.alignment = PP_ALIGN.CENTER

    # Standards badge
    badge = slide.shapes.add_textbox(Inches(6), Inches(5.5), Inches(3.5), Inches(0.8))
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = f"NGSS: {NGSS_STANDARD}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = MID_BLUE
    p.alignment = PP_ALIGN.RIGHT

    p = tf.add_paragraph()
    p.text = GRADE_LEVEL
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.RIGHT

    add_slide_number(slide, 1)

    # ========== SLIDE 2: LEARNING OBJECTIVES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "What You Will Learn Today"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Left column - Learning Goals
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Learning Goals"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    objectives = [
        "Identify the parts of California's fire system",
        "Explain how Earth's systems (water, air, land, life) interact",
        "Build a model showing cause-and-effect relationships",
        "Use simulations to predict what happens during drought"
    ]
    for obj in objectives:
        p = tf.add_paragraph()
        p.text = f"  *  {obj}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)

    # Right column - Key Vocabulary
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Vocabulary"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    vocab = [
        ("Component", "a part of a system"),
        ("Relationship", "how parts affect each other"),
        ("Simulation", "running a model to see what happens"),
        ("Earth System", "geosphere, hydrosphere, biosphere, atmosphere")
    ]
    for term, definition in vocab:
        p = tf.add_paragraph()
        p.text = f"  {term}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = f"     {definition}"
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = DARK_TEXT

    add_slide_number(slide, 2)

    # ========== SLIDE 3: THE QUESTION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Big Question"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    question = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1))
    tf = question.text_frame
    p = tf.paragraphs[0]
    p.text = "Why does California burn at the same time every year?"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Context box
    context = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(4.5), Inches(1.5))
    tf = context.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Every fall, California faces massive wildfires. But why fall? Why not spring or winter? Today we'll build a MODEL to discover the answer!"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT

    # Add landscape image with proper aspect ratio
    if images.get('landscape'):
        add_image_to_slide(slide, images['landscape'], 5.3, 3.2, 4.2, 3.2, center=True)

    add_slide_number(slide, 3)

    # ========== SLIDE 4: WHAT WE'LL BUILD ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Today We Will Build a Model!"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    content = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "A model helps us understand complex systems!"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = MID_BLUE

    steps = [
        ("1. LOCATE", "Identify the COMPONENTS (parts) of the fire system"),
        ("2. ESTABLISH", "Connect them with RELATIONSHIPS (+ or -)"),
        ("3. VISUALIZE", "Build your model in ModelIt!"),
        ("4. EVALUATE", "Run SIMULATIONS to test scenarios"),
        ("5. REVISE", "Improve your model based on evidence")
    ]
    for step, desc in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = BRAND_BLUE
        p.space_before = Pt(10)

        p = tf.add_paragraph()
        p.text = f"     {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT

    # Add modeling image with proper aspect ratio
    if images.get('modeling'):
        add_image_to_slide(slide, images['modeling'], 5.8, 2.5, 3.8, 3.5, center=True)

    add_slide_number(slide, 4)

    # ========== SLIDE 5: ACTIVITY 1 - SORT COMPONENTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Activity 1: Sort the Components"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Instructions
    instruct = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4.8), Inches(1))
    tf = instruct.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sort these components into EXTERNAL (inputs from outside) and INTERNAL (inside the system):"
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_TEXT

    # Components list
    comp_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(4.8), Inches(2.5))
    tf = comp_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Your Components:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    components = ["Rainfall", "Dry Vegetation", "Wind", "Fire Spread"]
    for comp in components:
        p = tf.add_paragraph()
        p.text = f"     *  {comp}"
        p.font.size = Pt(16)
        p.space_before = Pt(6)

    # Think-pair-share prompt
    think = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(4.8), Inches(1))
    tf = think.text_frame
    p = tf.paragraphs[0]
    p.text = "Think: Which components can we control? Which happen on their own?"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = MID_BLUE

    # Platform screenshot placeholder
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.6), Inches(2.1), Inches(4), Inches(3.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFA)
    box.line.color.rgb = MID_BLUE

    text_box = slide.shapes.add_textbox(Inches(5.8), Inches(3.7), Inches(3.6), Inches(1))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[ModelIt Platform Screenshot - Sorting Activity]"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 5)

    # ========== SLIDE 6: ACTIVITY 2 - CONNECT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Activity 2: Connect with Arrows"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(3))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Draw arrows to show HOW components affect each other:"
    p.font.size = Pt(17)

    p = tf.add_paragraph()
    p.text = "(+) POSITIVE Relationship"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x22, 0x8B, 0x22)  # Green
    p.space_before = Pt(14)

    p = tf.add_paragraph()
    p.text = "     When one goes UP, the other goes UP too"
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "(-) NEGATIVE Relationship"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xDC, 0x14, 0x3C)  # Red
    p.space_before = Pt(14)

    p = tf.add_paragraph()
    p.text = "     When one goes UP, the other goes DOWN"
    p.font.size = Pt(14)

    # Example question
    example = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5.5), Inches(1.5))
    tf = example.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Think About It:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    p = tf.add_paragraph()
    p.text = "When RAINFALL increases, does DRY VEGETATION increase or decrease?"
    p.font.size = Pt(15)
    p.font.italic = True
    p.space_before = Pt(6)

    # Add discussion image with proper aspect ratio
    if images.get('discussion'):
        add_image_to_slide(slide, images['discussion'], 5.8, 2.5, 3.8, 3.2, center=True)

    add_slide_number(slide, 6)

    # ========== SLIDE 7: ACTIVITY 3 - SIMULATE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Activity 3: Run the Simulation!"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5), Inches(4))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Test these scenarios in ModelIt!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    scenarios = [
        ("Scenario 1: DROUGHT", "Turn RAINFALL to OFF"),
        ("Scenario 2: WINDY DAY", "Turn WIND to ON"),
        ("Scenario 3: FIRE SEASON", "RAINFALL OFF + WIND ON")
    ]
    for name, instruction in scenarios:
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(12)

        p = tf.add_paragraph()
        p.text = f"     {instruction}"
        p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "\nWatch the activity graphs change!"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = MID_BLUE
    p.space_before = Pt(16)

    # Platform screenshot placeholder
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.3), Inches(2.1), Inches(4.3), Inches(4.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFA)
    box.line.color.rgb = MID_BLUE

    text_box = slide.shapes.add_textbox(Inches(5.5), Inches(4), Inches(3.9), Inches(1))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[ModelIt Platform Screenshot - Simulation Results Graph]"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 7)

    # ========== SLIDE 8: DISCUSSION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "What Did We Discover?"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    content = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.5), Inches(4.5))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Our Model Showed Us:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE

    discoveries = [
        "California's dry season = more dry vegetation = more fire risk",
        "Drought makes vegetation into FUEL for fires",
        "Wind helps fire SPREAD faster and farther",
        "These Earth systems work TOGETHER to create fire conditions"
    ]
    for d in discoveries:
        p = tf.add_paragraph()
        p.text = f"  *  {d}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(10)

    # Answer box
    answer = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.2))
    tf = answer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Answer: California burns in fall because summer drought dries out vegetation (FUEL), and fall brings Santa Ana winds that spread fire rapidly!"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = MID_BLUE

    # Add cover image again for discussion
    if images.get('cover'):
        add_image_to_slide(slide, images['cover'], 6.2, 2.5, 3.3, 3, center=True)

    add_slide_number(slide, 8)

    # ========== SLIDE 9: STEM CHALLENGE PREVIEW ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "STEM Challenge: Firebreak Engineers!"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    content = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "YOUR ENGINEERING MISSION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE

    p = tf.add_paragraph()
    p.text = "You learned that DRY VEGETATION is FUEL. Real engineers create FIREBREAKS - areas with no fuel - to stop fires!"
    p.font.size = Pt(14)
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = "\nThe Challenge:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = "Design firebreaks to protect Willow Creek from wildfire. You have a budget for only 3 firebreaks."
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = "\nThink Like an Engineer:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BRAND_BLUE
    p.space_before = Pt(10)

    engineer_qs = [
        "Which direction does wind blow? (From the east!)",
        "Does fire spread uphill or downhill faster?",
        "Where should firebreaks go to protect the town?"
    ]
    for q in engineer_qs:
        p = tf.add_paragraph()
        p.text = f"     *  {q}"
        p.font.size = Pt(13)
        p.space_before = Pt(4)

    # Add STEM image with proper aspect ratio
    if images.get('stem'):
        add_image_to_slide(slide, images['stem'], 5.8, 2.5, 3.8, 3.5, center=True)

    # Career Connection blurb
    career_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9.0), Inches(0.95))
    career_box.fill.solid()
    career_box.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)  # NAVY background
    tf = career_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)

    p = tf.paragraphs[0]
    p.text = "REAL CAREER CONNECTION: "
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xA5, 0x00)  # orange label

    run = p.add_run()
    run.text = "Fire Protection Engineers design real firebreaks to protect communities — and earn $95,000–$150,000/year. "
    run.font.size = Pt(11)
    run.font.bold = False
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    run2 = p.add_run()
    run2.text = "Environmental Scientists studying wildfire systems earn $75,000–$120,000/year. The skills you're using TODAY are the same ones they use on the job!"
    run2.font.size = Pt(11)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_slide_number(slide, 9)

    prs.save(output_path)
    print(f"[OK] Created: {output_path}")
    return output_path


# ============================================
# WORD DOCUMENT HELPERS
# ============================================

def set_cell_shading(cell, color_hex):
    """Set background color of a table cell"""
    shading = OxmlElement('w:shd')
    shading.set(qn('w:fill'), color_hex)
    cell._tc.get_or_add_tcPr().append(shading)


def add_name_date_line(doc):
    """Add Name/Date line to page"""
    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    row = table.rows[0]
    row.cells[0].text = "Name: _______________________________"
    row.cells[1].text = "Date: _______________"
    for cell in row.cells:
        cell.paragraphs[0].runs[0].font.size = DocxPt(11)
    doc.add_paragraph()


def add_section_header(doc, text, level=1):
    """Add a styled section header"""
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = True
    if level == 1:
        run.font.size = DocxPt(24)
        run.font.color.rgb = DOCX_NAVY
    elif level == 2:
        run.font.size = DocxPt(18)
        run.font.color.rgb = DOCX_BRAND_BLUE
    elif level == 3:
        run.font.size = DocxPt(14)
        run.font.color.rgb = DOCX_MID_BLUE
    return p


def add_response_box(doc, height_lines=3, prompt=""):
    """Add a response box for student answers"""
    if prompt:
        p = doc.add_paragraph()
        p.add_run(prompt).font.size = DocxPt(11)

    table = doc.add_table(rows=1, cols=1)
    table.style = 'Table Grid'
    cell = table.rows[0].cells[0]
    cell.text = "\n" * height_lines
    doc.add_paragraph()
    return table


# ============================================
# STUDENT ACTIVITY PACK (WORD) - ENHANCED
# ============================================

def create_student_activity_pack():
    """Create the enhanced student activity pack with open response questions."""
    output_path = os.path.join(OUTPUT_DIR, f'{LESSON_ID}-Student-Activity-Pack.docx')
    doc = Document()

    # Set margins
    for section in doc.sections:
        section.top_margin = DocxInches(0.7)
        section.bottom_margin = DocxInches(0.7)
        section.left_margin = DocxInches(0.7)
        section.right_margin = DocxInches(0.7)

    # ========== COVER PAGE ==========
    doc.add_paragraph()
    doc.add_paragraph()

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("ModelIt!")
    run.bold = True
    run.font.size = DocxPt(48)
    run.font.color.rgb = DOCX_NAVY

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("Student Activity Pack")
    run.font.size = DocxPt(28)
    run.font.color.rgb = DOCX_BRAND_BLUE

    doc.add_paragraph()

    lesson_title = doc.add_paragraph()
    lesson_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = lesson_title.add_run(LESSON_TITLE)
    run.bold = True
    run.font.size = DocxPt(30)
    run.font.color.rgb = DOCX_NAVY

    lesson_sub = doc.add_paragraph()
    lesson_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = lesson_sub.add_run(LESSON_SUBTITLE)
    run.italic = True
    run.font.size = DocxPt(15)

    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()

    # Name/Date on cover
    add_name_date_line(doc)

    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = info.add_run(f"Grade Level: {GRADE_LEVEL}  |  Standard: {NGSS_STANDARD}")
    run.font.size = DocxPt(12)

    doc.add_page_break()

    # ========== TABLE OF CONTENTS ==========
    add_section_header(doc, "What's Inside")
    add_name_date_line(doc)

    toc_items = [
        ("1. Pre-Assessment", "What do you already know?"),
        ("2. Scoring Rubric", "How your work will be assessed"),
        ("3. Component Analysis", "Sort and classify system parts"),
        ("4. Model Recording Page", "Sketch your model"),
        ("5. Simulation Observations", "Record what happens"),
        ("6. Research & Extend", "Add new components"),
        ("7. Reflection Questions", "What did you learn?"),
        ("8. STEM Challenge", "Firebreak Engineers!"),
    ]

    for item, desc in toc_items:
        p = doc.add_paragraph()
        run = p.add_run(item)
        run.bold = True
        run.font.size = DocxPt(13)
        p.add_run(f"  -  {desc}").font.size = DocxPt(11)

    doc.add_page_break()

    # ========== PRE-ASSESSMENT ==========
    add_section_header(doc, "Pre-Assessment: What Do You Know?")
    add_name_date_line(doc)

    p = doc.add_paragraph()
    p.add_run("Answer these questions BEFORE the lesson begins.").font.size = DocxPt(11)

    doc.add_paragraph()

    # Open response questions
    questions = [
        "1. What do you already know about wildfires? Have you seen news about California fires?",
        "2. What do you think causes wildfires to spread so quickly?",
        "3. What is a model? Why do scientists use models?",
        "4. What is a system? Can you name some examples of systems?"
    ]

    for q in questions:
        p = doc.add_paragraph()
        run = p.add_run(q)
        run.bold = True
        run.font.size = DocxPt(11)
        add_response_box(doc, 3)

    doc.add_page_break()

    # ========== SCORING RUBRIC ==========
    add_section_header(doc, "Scoring Rubric")
    add_name_date_line(doc)

    p = doc.add_paragraph()
    p.add_run("This rubric shows how your modeling work will be assessed. Use it to check your own work!").font.size = DocxPt(11)

    doc.add_paragraph()

    rubric_data = [
        ["Skill", "4 - Expert", "3 - Proficient", "2 - Developing", "1 - Beginning"],
        ["Identifying\nComponents", "Identifies all components correctly and explains their role", "Identifies most components correctly", "Identifies some components; may misclassify", "Struggles to identify components"],
        ["Establishing\nRelationships", "All arrows correct direction with +/- signs and can explain why", "Most relationships correct; minor errors", "Some relationships correct; direction or sign errors", "Relationships incorrect or missing"],
        ["Running\nSimulations", "Tests multiple scenarios systematically and predicts outcomes", "Runs all scenarios and records observations", "Runs some scenarios with guidance", "Needs significant help to run simulation"],
        ["Explaining\nResults", "Clear explanation with scientific evidence from model", "Good explanation connecting to model", "Partial explanation; missing connections", "Struggles to explain results"],
    ]

    table = doc.add_table(rows=len(rubric_data), cols=5)
    table.style = 'Table Grid'

    for i, row_data in enumerate(rubric_data):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = cell_text
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = DocxPt(9)
            if i == 0:
                cell.paragraphs[0].runs[0].bold = True
                set_cell_shading(cell, "1A4780")
                cell.paragraphs[0].runs[0].font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)

    doc.add_page_break()

    # ========== COMPONENT ANALYSIS ==========
    add_section_header(doc, "Component Analysis")
    add_name_date_line(doc)

    add_section_header(doc, "What is a Component?", level=2)
    p = doc.add_paragraph()
    p.add_run("A ").font.size = DocxPt(11)
    run = p.add_run("component")
    run.bold = True
    run.font.size = DocxPt(11)
    p.add_run(" is a part of a system. In our California wildfire system, we have these components:").font.size = DocxPt(11)

    doc.add_paragraph()

    # Components list with descriptions
    components = [
        ("Rainfall", "Water that falls from clouds; affects how wet or dry plants are"),
        ("Dry Vegetation", "Dead grass, leaves, and dry plants that can burn easily"),
        ("Wind", "Moving air that can push fire and carry embers"),
        ("Fire Spread", "How quickly and far fire moves across the land")
    ]

    for name, desc in components:
        p = doc.add_paragraph()
        run = p.add_run(f"  *  {name}: ")
        run.bold = True
        run.font.size = DocxPt(11)
        p.add_run(desc).font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "Sorting Activity", level=2)
    p = doc.add_paragraph()
    p.add_run("Sort the components into EXTERNAL (inputs from outside the system) and INTERNAL (processes inside the system):").font.size = DocxPt(11)

    doc.add_paragraph()

    comp_table = doc.add_table(rows=2, cols=2)
    comp_table.style = 'Table Grid'
    labels = ["EXTERNAL Components\n(Come from outside)", "INTERNAL Components\n(Happen inside system)"]
    for i, cell in enumerate(comp_table.rows[0].cells):
        cell.text = labels[i]
        cell.paragraphs[0].runs[0].bold = True
        set_cell_shading(cell, "7EC8E3")
    for cell in comp_table.rows[1].cells:
        cell.text = "\n\n\n\n"

    doc.add_paragraph()

    # Open response question
    p = doc.add_paragraph()
    run = p.add_run("Explain your reasoning: Why did you classify the components this way?")
    run.bold = True
    run.font.size = DocxPt(11)
    add_response_box(doc, 4)

    doc.add_page_break()

    # ========== MODEL RECORDING PAGE ==========
    add_section_header(doc, "Model Recording Page")
    add_name_date_line(doc)

    p = doc.add_paragraph()
    p.add_run("Sketch your model showing all components and their relationships.").font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "Relationship Key:", level=3)
    key_items = [
        ("(+) POSITIVE:", "When one component increases, the other ALSO increases"),
        ("(-) NEGATIVE:", "When one component increases, the other DECREASES")
    ]
    for label, desc in key_items:
        p = doc.add_paragraph()
        run = p.add_run(label + " ")
        run.bold = True
        run.font.size = DocxPt(11)
        p.add_run(desc).font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "My Model Sketch:", level=3)
    p = doc.add_paragraph()
    p.add_run("Draw your components as boxes/circles. Connect them with arrows showing direction. Label each arrow as + or -.").font.size = DocxPt(10)
    p.italic = True

    sketch_table = doc.add_table(rows=1, cols=1)
    sketch_table.style = 'Table Grid'
    sketch_table.rows[0].cells[0].text = "\n\n\n\n\n\n\n\n\n\n\n"

    doc.add_paragraph()

    # Open response questions about relationships
    p = doc.add_paragraph()
    run = p.add_run("Explain ONE relationship in your model. Why does it work that way?")
    run.bold = True
    run.font.size = DocxPt(11)
    add_response_box(doc, 4)

    doc.add_page_break()

    # ========== SIMULATION OBSERVATIONS ==========
    add_section_header(doc, "Simulation Observations")
    add_name_date_line(doc)

    p = doc.add_paragraph()
    p.add_run("Run each scenario in ModelIt! and record your observations. Pay attention to how the activity graph changes!").font.size = DocxPt(11)

    doc.add_paragraph()

    scenarios = [
        ("Scenario 1: DROUGHT", "Turn RAINFALL to OFF. Leave WIND OFF.",
         "Prediction: What do you think will happen to DRY VEGETATION and FIRE SPREAD?"),
        ("Scenario 2: WINDY DAY", "Turn WIND to ON. Leave RAINFALL ON.",
         "Prediction: What do you think will happen to FIRE SPREAD?"),
        ("Scenario 3: CALIFORNIA FIRE SEASON", "Turn RAINFALL OFF and WIND ON.",
         "Prediction: This is what happens every fall in California. What do you expect?"),
    ]

    for title, instruction, prediction in scenarios:
        add_section_header(doc, title, level=2)

        p = doc.add_paragraph()
        run = p.add_run("Settings: ")
        run.bold = True
        run.font.size = DocxPt(11)
        p.add_run(instruction).font.size = DocxPt(11)

        # Prediction
        p = doc.add_paragraph()
        run = p.add_run(prediction)
        run.font.size = DocxPt(11)
        run.italic = True
        add_response_box(doc, 2)

        # Observation
        p = doc.add_paragraph()
        run = p.add_run("Observation: What actually happened? Describe what you saw on the graph.")
        run.bold = True
        run.font.size = DocxPt(11)
        add_response_box(doc, 3)

    doc.add_page_break()

    # ========== RESEARCH & EXTEND ==========
    add_section_header(doc, "Research & Extend")
    add_name_date_line(doc)

    p = doc.add_paragraph()
    p.add_run("Real systems are more complex than our basic model. Use what you've learned to EXTEND your model with new components!").font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "Possible New Components:", level=2)
    examples = [
        ("Temperature", "How might hot weather affect dry vegetation and fire?"),
        ("Humidity", "How does moisture in the air affect fire spread?"),
        ("Terrain/Slopes", "Does fire spread faster uphill or downhill?"),
        ("Santa Ana Winds", "These are special dry, hot winds from the east. How might they be different from regular wind?"),
        ("Human Activity", "How might people affect the fire system? (Both causing and preventing fires)")
    ]
    for ex, explanation in examples:
        p = doc.add_paragraph()
        run = p.add_run(f"  *  {ex}: ")
        run.bold = True
        run.font.size = DocxPt(11)
        p.add_run(explanation).font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "My Extension:", level=2)

    p = doc.add_paragraph()
    run = p.add_run("New component I want to add: ")
    run.bold = True
    run.font.size = DocxPt(11)
    p.add_run("________________________________").font.size = DocxPt(11)

    p = doc.add_paragraph()
    run = p.add_run("How does it connect to existing components? (Draw and explain)")
    run.bold = True
    run.font.size = DocxPt(11)

    sketch = doc.add_table(rows=1, cols=1)
    sketch.style = 'Table Grid'
    sketch.rows[0].cells[0].text = "\n\n\n\n\n"

    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run("What evidence or reasoning supports these connections?")
    run.bold = True
    run.font.size = DocxPt(11)
    add_response_box(doc, 4)

    doc.add_page_break()

    # ========== REFLECTION QUESTIONS ==========
    add_section_header(doc, "Reflection Questions")
    add_name_date_line(doc)

    p = doc.add_paragraph()
    p.add_run("Answer these questions after completing all activities.").font.size = DocxPt(11)

    doc.add_paragraph()

    reflections = [
        "1. Answer the Big Question: Why does California burn at the same time every year? Use evidence from your model!",
        "2. How do Earth's systems (geosphere, hydrosphere, biosphere, atmosphere) interact to create fire conditions?",
        "3. What was the most surprising thing you discovered from running the simulations?",
        "4. How could you use what you learned to help communities prepare for fire season?",
        "5. What questions do you still have about wildfires or modeling systems?"
    ]

    for q in reflections:
        p = doc.add_paragraph()
        run = p.add_run(q)
        run.bold = True
        run.font.size = DocxPt(11)
        add_response_box(doc, 4)

    doc.add_page_break()

    # ========== STEM CHALLENGE ==========
    add_section_header(doc, "STEM CHALLENGE: Firebreak Engineers")
    add_name_date_line(doc)

    add_section_header(doc, "YOUR MISSION", level=2)
    p = doc.add_paragraph()
    p.add_run("You learned that DRY VEGETATION is the FUEL for fire spread. Real engineers use FIREBREAKS - areas with no fuel - to stop fires! Now YOU will design firebreaks to protect a community.").font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "THE SCENARIO", level=2)
    p = doc.add_paragraph()
    p.add_run("Willow Creek is a small California town surrounded by forest. Fire season is coming! Santa Ana winds blow from the EAST. The town has a budget for only 3 firebreaks. Where will you put them?").font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "ENGINEERING DESIGN THINKING", level=2)
    questions = [
        "Which direction does fire spread fastest? (Think about wind direction!)",
        "Does fire spread faster uphill or downhill? Why?",
        "Should firebreaks be between the fire source and the town, or behind the town?",
        "How does this challenge connect to what you learned in your model?"
    ]
    for q in questions:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(q).font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "MY FIREBREAK DESIGN", level=2)
    p = doc.add_paragraph()
    p.add_run("Draw a map showing: Town center, Forest around it, Wind direction arrows, Your 3 firebreaks").font.size = DocxPt(10)
    p.italic = True

    design_box = doc.add_table(rows=1, cols=1)
    design_box.style = 'Table Grid'
    design_box.rows[0].cells[0].text = "\n\n\n\n\n\n\n\n\n\n"

    doc.add_paragraph()

    add_section_header(doc, "EXPLAIN YOUR DESIGN", level=2)
    p = doc.add_paragraph()
    run = p.add_run("Why did you put the firebreaks where you did? Use scientific reasoning from your model!")
    run.bold = True
    run.font.size = DocxPt(11)
    add_response_box(doc, 5)

    doc.save(output_path)
    print(f"[OK] Created: {output_path}")
    return output_path


# ============================================
# TEACHER'S GUIDE (WORD) - COMPREHENSIVE
# ============================================

def create_teachers_guide():
    """Create comprehensive teacher's guide with CAST info, facilitation guide, and expected responses."""
    output_path = os.path.join(OUTPUT_DIR, f'{LESSON_ID}-Teachers-Guide.docx')
    doc = Document()

    for section in doc.sections:
        section.top_margin = DocxInches(0.6)
        section.bottom_margin = DocxInches(0.6)
        section.left_margin = DocxInches(0.6)
        section.right_margin = DocxInches(0.6)

    # ========== COVER ==========
    doc.add_paragraph()
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("ModelIt! Teacher's Guide")
    run.bold = True
    run.font.size = DocxPt(36)
    run.font.color.rgb = DOCX_NAVY

    doc.add_paragraph()
    lesson = doc.add_paragraph()
    lesson.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = lesson.add_run(f"{LESSON_ID}: {LESSON_TITLE}")
    run.bold = True
    run.font.size = DocxPt(22)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(LESSON_SUBTITLE)
    run.font.size = DocxPt(14)
    run.italic = True

    doc.add_paragraph()

    # Quick reference table
    quick_ref = [
        ["Grade Level", GRADE_LEVEL],
        ["NGSS Standard", f"{NGSS_STANDARD}: Develop a model using an example to describe ways the geosphere, biosphere, hydrosphere, and/or atmosphere interact."],
        ["Time Required", "45-60 minutes (can extend to 90 min with STEM challenge)"],
        ["Materials", "Student devices (Chromebooks/tablets), Activity Pack, Internet access"],
        ["Prep Required", "Load ModelIt! platform; print Activity Packs; review lesson"],
        ["Grouping", "Pairs or small groups of 3 (recommended)"],
    ]

    table = doc.add_table(rows=len(quick_ref), cols=2)
    table.style = 'Table Grid'
    for i, (label, value) in enumerate(quick_ref):
        table.rows[i].cells[0].text = label
        for para in table.rows[i].cells[0].paragraphs:
            for run in para.runs:
                run.bold = True
                run.font.size = DocxPt(10)
        set_cell_shading(table.rows[i].cells[0], "7EC8E3")
        table.rows[i].cells[1].text = value
        for para in table.rows[i].cells[1].paragraphs:
            for run in para.runs:
                run.font.size = DocxPt(10)

    doc.add_page_break()

    # ========== TABLE OF CONTENTS ==========
    add_section_header(doc, "Guide Contents")

    toc_items = [
        "1. NGSS Standards Unpacking & CAST Alignment",
        "2. Background Content Knowledge for Teachers",
        "3. The LEVER Framework",
        "4. Slide-by-Slide Facilitation Guide",
        "5. Answer Key with Expected Student Responses",
        "6. STEM Challenge Teacher Guidance",
        "7. Differentiation & Extensions",
        "8. Common Misconceptions & How to Address Them"
    ]

    for item in toc_items:
        p = doc.add_paragraph()
        p.add_run(item).font.size = DocxPt(12)

    doc.add_page_break()

    # ========== 1. NGSS STANDARDS UNPACKING & CAST ALIGNMENT ==========
    add_section_header(doc, "1. NGSS Standards Unpacking & CAST Alignment")

    add_section_header(doc, "Performance Expectation: 5-ESS2-1", level=2)
    p = doc.add_paragraph()
    run = p.add_run("Develop a model using an example to describe ways the geosphere, biosphere, hydrosphere, and/or atmosphere interact.")
    run.italic = True
    run.font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "Three-Dimensional Learning", level=2)

    dimensions = [
        ("Science & Engineering Practice", "Developing and Using Models",
         "Students build a computational model to represent the California wildfire system, identifying components and their relationships."),
        ("Disciplinary Core Idea", "ESS2.A: Earth Materials and Systems",
         "Earth's major systems interact in complex ways. The fire system demonstrates interactions between the hydrosphere (rainfall), atmosphere (wind), biosphere (vegetation), and geosphere (terrain)."),
        ("Crosscutting Concept", "Systems and System Models",
         "Students learn that models can represent systems and their interactions. They analyze cause-and-effect relationships within the system.")
    ]

    for dim_name, dim_title, description in dimensions:
        p = doc.add_paragraph()
        run = p.add_run(dim_name + ": ")
        run.bold = True
        run.font.size = DocxPt(11)
        run = p.add_run(dim_title)
        run.italic = True
        run.font.size = DocxPt(11)

        p = doc.add_paragraph()
        p.add_run(description).font.size = DocxPt(11)
        doc.add_paragraph()

    add_section_header(doc, "CAST Testing Alignment", level=2)

    p = doc.add_paragraph()
    p.add_run("The California Science Test (CAST) assesses 5-ESS2-1 through performance tasks that ask students to:").font.size = DocxPt(11)

    cast_items = [
        "Identify components of Earth systems from diagrams or descriptions",
        "Analyze cause-and-effect relationships between Earth systems",
        "Use models to explain or predict phenomena involving multiple Earth systems",
        "Evaluate how changes in one system affect other systems"
    ]
    for item in cast_items:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(item).font.size = DocxPt(11)

    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run("Sample CAST-Style Questions Students Should Be Able to Answer:")
    run.bold = True
    run.font.size = DocxPt(11)

    cast_questions = [
        ("Multiple Choice:", "Which Earth system provides the fuel for wildfires? A) Hydrosphere B) Biosphere C) Atmosphere D) Geosphere [Answer: B - vegetation is part of the biosphere]"),
        ("Constructed Response:", "Explain how a decrease in rainfall affects the other components of the California fire system. Use cause-and-effect language in your response."),
        ("Model Interpretation:", "Look at this diagram of a fire system. Draw an arrow showing how wind affects fire spread. Label it as positive (+) or negative (-). Explain your reasoning.")
    ]

    for q_type, q_text in cast_questions:
        p = doc.add_paragraph()
        run = p.add_run(q_type + " ")
        run.bold = True
        run.font.size = DocxPt(10)
        p.add_run(q_text).font.size = DocxPt(10)

    doc.add_page_break()

    # ========== 2. BACKGROUND CONTENT KNOWLEDGE ==========
    add_section_header(doc, "2. Background Content Knowledge for Teachers")

    add_section_header(doc, "California Wildfire Science", level=2)

    p = doc.add_paragraph()
    p.add_run("California experiences predictable annual fire seasons due to a confluence of Earth system interactions:").font.size = DocxPt(11)

    doc.add_paragraph()

    knowledge_sections = [
        ("The Mediterranean Climate Cycle",
         "California has a Mediterranean climate with wet winters and dry summers. From May through October, most areas receive virtually no rainfall. This extended dry period causes vegetation to lose moisture content, dropping from ~50% to as low as 5% moisture - essentially turning living plants into kindling."),
        ("Vegetation as Fuel",
         "The biosphere component is critical. Chaparral, grasslands, and forest understory accumulate over multiple years. When moisture content drops below 15%, vegetation becomes highly flammable. Invasive grasses (like cheatgrass) dry out faster than native plants, increasing fire risk."),
        ("Santa Ana Winds",
         "These are katabatic winds that blow from the Great Basin deserts toward the coast, typically from October through February. They are hot, dry (sometimes <10% humidity), and can reach 40-70 mph. The combination of timing (when vegetation is driest) and conditions (hot, dry, windy) creates extreme fire weather."),
        ("Fire Behavior Physics",
         "Fire spreads through heat transfer: convection (hot gases rise), radiation (heat radiates outward), and direct flame contact. Wind increases all three mechanisms and pushes flames forward. Fire also spreads faster uphill because flames preheat fuel above them through convection.")
    ]

    for section_title, content in knowledge_sections:
        p = doc.add_paragraph()
        run = p.add_run(section_title + ": ")
        run.bold = True
        run.font.size = DocxPt(11)
        run.font.color.rgb = DOCX_BRAND_BLUE

        p = doc.add_paragraph()
        p.add_run(content).font.size = DocxPt(11)
        doc.add_paragraph()

    add_section_header(doc, "Earth Systems Connections", level=2)

    systems_table = [
        ["Earth System", "Role in Wildfire", "Example"],
        ["Hydrosphere", "Controls vegetation moisture", "Less rain = drier vegetation = more fuel"],
        ["Biosphere", "Provides fuel (vegetation)", "Plants die and dry out, becoming flammable"],
        ["Atmosphere", "Provides oxygen; wind spreads fire", "Santa Ana winds push fire rapidly"],
        ["Geosphere", "Terrain affects fire spread", "Fire moves faster uphill; valleys channel wind"]
    ]

    sys_table = doc.add_table(rows=len(systems_table), cols=3)
    sys_table.style = 'Table Grid'
    for i, row_data in enumerate(systems_table):
        for j, cell_text in enumerate(row_data):
            cell = sys_table.rows[i].cells[j]
            cell.text = cell_text
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = DocxPt(10)
            if i == 0:
                cell.paragraphs[0].runs[0].bold = True
                set_cell_shading(cell, "1A4780")
                cell.paragraphs[0].runs[0].font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)

    doc.add_page_break()

    # ========== 3. THE LEVER FRAMEWORK ==========
    add_section_header(doc, "3. The LEVER Framework")

    p = doc.add_paragraph()
    p.add_run("LEVER is the instructional framework that guides students through authentic scientific modeling. Each letter represents a phase of model development:").font.size = DocxPt(11)

    doc.add_paragraph()

    lever_phases = [
        ("L - LOCATE the System",
         "Students identify the components (parts) of the system they're studying.",
         "Students identify: Rainfall, Dry Vegetation, Wind, Fire Spread as the key components of California's fire system."),
        ("E - ESTABLISH Relationships",
         "Students determine how components affect each other using positive (+) and negative (-) relationships.",
         "Students connect Rainfall --> Dry Vegetation with (-) because more rain = less dry vegetation. They connect Dry Vegetation --> Fire Spread with (+) because more fuel = more fire."),
        ("V - VISUALIZE & Model",
         "Students build their model in the ModelIt! platform, arranging components and drawing connections.",
         "Students drag components onto the canvas, sort them into External/Internal, and draw directional arrows."),
        ("E - EVALUATE Outcomes",
         "Students run simulations and observe what happens under different conditions.",
         "Students test drought conditions (Rainfall OFF) and observe how activity graphs show Dry Vegetation and Fire Spread increasing."),
        ("R - REVISE & Extend",
         "Students improve their model based on evidence, adding new components from research.",
         "Students might add Temperature, Humidity, or Santa Ana Winds based on readings about fire science.")
    ]

    for phase, description, example in lever_phases:
        p = doc.add_paragraph()
        run = p.add_run(phase)
        run.bold = True
        run.font.size = DocxPt(12)
        run.font.color.rgb = DOCX_BRAND_BLUE

        p = doc.add_paragraph()
        p.add_run(description).font.size = DocxPt(11)

        p = doc.add_paragraph()
        run = p.add_run("In this lesson: ")
        run.bold = True
        run.font.size = DocxPt(10)
        run.font.color.rgb = DOCX_MID_BLUE
        p.add_run(example).font.size = DocxPt(10)
        doc.add_paragraph()

    doc.add_page_break()

    # ========== 4. SLIDE-BY-SLIDE FACILITATION GUIDE ==========
    add_section_header(doc, "4. Slide-by-Slide Facilitation Guide")

    p = doc.add_paragraph()
    p.add_run("Use this guide alongside the PowerPoint. Each slide includes what students see, what you say, and what to watch for.").font.size = DocxPt(11)

    doc.add_paragraph()

    slides = [
        {
            "number": "Slide 1",
            "title": "Cover Slide",
            "visual": "Title: 'When Trees Become Matches'\nSubtitle about California's burning season\nModelIt! logo\nImage of students collaborating",
            "say": "\"Today we're going to investigate a California mystery - why does our state burn at the same time every single year? We'll use a tool called ModelIt! to build a model that helps us discover the answer.\"",
            "do": "Build curiosity! Ask: 'Has anyone seen news about California fires? When do they usually happen?' Accept all responses - they'll confirm their ideas through the model.",
            "time": "2 min"
        },
        {
            "number": "Slide 2",
            "title": "Learning Objectives",
            "visual": "Two columns:\n- Learning Goals (4 objectives)\n- Key Vocabulary (4 terms)",
            "say": "\"Here's what we'll accomplish today. Notice these vocabulary words - we'll use them as we build our model. A component is a part of a system. A relationship shows how parts affect each other.\"",
            "do": "Don't spend too long here - these will be reinforced throughout. Quick check: 'Who can give me an example of a system?' (car, body, weather)",
            "time": "1-2 min"
        },
        {
            "number": "Slide 3",
            "title": "The Big Question",
            "visual": "Question: 'Why does California burn at the same time every year?'\nImage of California hillside during fire season",
            "say": "\"This is our driving question. Every fall, California faces massive wildfires. But why fall specifically? Why not spring or winter? By the end of today, your model will help you answer this!\"",
            "do": "Have students make predictions in their Activity Pack (Pre-Assessment). Collect a few ideas verbally but don't correct them yet.",
            "time": "3 min"
        },
        {
            "number": "Slide 4",
            "title": "What We'll Build",
            "visual": "LEVER steps listed\nImage of student building model",
            "say": "\"We'll follow these five steps to build our model. LOCATE the parts, ESTABLISH relationships, VISUALIZE by building, EVALUATE by testing, and REVISE to improve. This is exactly what real scientists do!\"",
            "do": "Emphasize that modeling is a PRACTICE scientists use. Connect to students' experience: 'Have you ever built a model of something? What did it help you understand?'",
            "time": "2 min"
        },
        {
            "number": "Slide 5",
            "title": "Activity 1: Sort Components",
            "visual": "List of 4 components\nSorting instruction\nPlatform screenshot placeholder",
            "say": "\"Here are our four components. Your job is to sort them - which ones are EXTERNAL (inputs from outside the system, like weather) and which are INTERNAL (things happening inside the system)?\"",
            "do": "THIS IS HANDS-ON TIME. Have students open ModelIt! and their Activity Packs. Walk around as pairs discuss and sort. Listen for reasoning, not just answers.",
            "time": "5-7 min"
        },
        {
            "number": "Slide 6",
            "title": "Activity 2: Connect with Arrows",
            "visual": "(+) and (-) relationship explanation\nThink question about Rainfall --> Dry Vegetation",
            "say": "\"Now we connect our components with arrows. But arrows aren't just about WHAT connects - they're about HOW. Is the relationship positive (both go up together) or negative (one up, one down)?\"",
            "do": "Model the thinking aloud: 'If rainfall INCREASES, what happens to dry vegetation? It gets WET, so dry vegetation DECREASES. That's a NEGATIVE relationship!' Have pairs draw their arrows.",
            "time": "7-10 min"
        },
        {
            "number": "Slide 7",
            "title": "Activity 3: Run Simulation",
            "visual": "Three scenarios to test\nPlatform screenshot placeholder",
            "say": "\"Now comes the fun part - let's see if your model actually works! Test each scenario and watch what happens to the activity graphs. Compare what you predicted to what actually happens.\"",
            "do": "Circulate as students run simulations. Ask probing questions: 'What happened when you turned rainfall off?' 'Why do you think fire spread increased?' Have students record in Activity Packs.",
            "time": "8-10 min"
        },
        {
            "number": "Slide 8",
            "title": "Discussion",
            "visual": "List of discoveries\nAnswer box explaining California fire timing",
            "say": "\"What did your simulations reveal? Let's put it together - California burns in fall because summer drought dries out vegetation, making fuel, and then fall Santa Ana winds spread fire rapidly!\"",
            "do": "Facilitate a class discussion. Have multiple groups share what they observed. Connect findings to the Big Question. Students should feel like THEY discovered the answer, not that you told them.",
            "time": "5 min"
        },
        {
            "number": "Slide 9",
            "title": "STEM Challenge Preview",
            "visual": "Firebreak Engineers mission\nEngineering thinking questions\nImage of students engineering",
            "say": "\"You discovered that dry vegetation is FUEL. Real engineers use this knowledge to create FIREBREAKS - areas with no fuel so fire can't spread. Now YOU will design firebreaks to protect a community!\"",
            "do": "If time allows, have students begin the STEM Challenge in their Activity Packs. This can also be homework or a Day 2 activity.",
            "time": "5-10 min"
        }
    ]

    for slide in slides:
        # Slide header
        p = doc.add_paragraph()
        run = p.add_run(f"{slide['number']}: {slide['title']}")
        run.bold = True
        run.font.size = DocxPt(14)
        run.font.color.rgb = DOCX_BRAND_BLUE

        # Create two-column table for visual and notes
        slide_table = doc.add_table(rows=4, cols=2)
        slide_table.style = 'Table Grid'

        # Row 1: Visual
        slide_table.rows[0].cells[0].text = "What Students See:"
        set_cell_shading(slide_table.rows[0].cells[0], "E8F4F8")
        slide_table.rows[0].cells[0].paragraphs[0].runs[0].bold = True
        slide_table.rows[0].cells[0].paragraphs[0].runs[0].font.size = DocxPt(9)
        slide_table.rows[0].cells[1].text = slide['visual']
        slide_table.rows[0].cells[1].paragraphs[0].runs[0].font.size = DocxPt(9)

        # Row 2: What to Say
        slide_table.rows[1].cells[0].text = "What to Say:"
        set_cell_shading(slide_table.rows[1].cells[0], "E8F4F8")
        slide_table.rows[1].cells[0].paragraphs[0].runs[0].bold = True
        slide_table.rows[1].cells[0].paragraphs[0].runs[0].font.size = DocxPt(9)
        slide_table.rows[1].cells[1].text = slide['say']
        slide_table.rows[1].cells[1].paragraphs[0].runs[0].font.size = DocxPt(9)

        # Row 3: What to Do
        slide_table.rows[2].cells[0].text = "What to Do/Watch For:"
        set_cell_shading(slide_table.rows[2].cells[0], "E8F4F8")
        slide_table.rows[2].cells[0].paragraphs[0].runs[0].bold = True
        slide_table.rows[2].cells[0].paragraphs[0].runs[0].font.size = DocxPt(9)
        slide_table.rows[2].cells[1].text = slide['do']
        slide_table.rows[2].cells[1].paragraphs[0].runs[0].font.size = DocxPt(9)

        # Row 4: Time
        slide_table.rows[3].cells[0].text = "Approximate Time:"
        set_cell_shading(slide_table.rows[3].cells[0], "E8F4F8")
        slide_table.rows[3].cells[0].paragraphs[0].runs[0].bold = True
        slide_table.rows[3].cells[0].paragraphs[0].runs[0].font.size = DocxPt(9)
        slide_table.rows[3].cells[1].text = slide['time']
        slide_table.rows[3].cells[1].paragraphs[0].runs[0].font.size = DocxPt(9)

        doc.add_paragraph()

    doc.add_page_break()

    # ========== 5. ANSWER KEY WITH EXPECTED RESPONSES ==========
    add_section_header(doc, "5. Answer Key with Expected Student Responses")

    add_section_header(doc, "Component Sorting", level=2)

    sorting_table = doc.add_table(rows=3, cols=2)
    sorting_table.style = 'Table Grid'
    sorting_table.rows[0].cells[0].text = "EXTERNAL Components"
    sorting_table.rows[0].cells[1].text = "INTERNAL Components"
    for cell in sorting_table.rows[0].cells:
        cell.paragraphs[0].runs[0].bold = True
        set_cell_shading(cell, "1A4780")
        cell.paragraphs[0].runs[0].font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)
        cell.paragraphs[0].runs[0].font.size = DocxPt(10)

    sorting_table.rows[1].cells[0].text = "Rainfall"
    sorting_table.rows[1].cells[1].text = "Dry Vegetation"
    sorting_table.rows[2].cells[0].text = "Wind"
    sorting_table.rows[2].cells[1].text = "Fire Spread"
    for row in sorting_table.rows[1:]:
        for cell in row.cells:
            cell.paragraphs[0].runs[0].font.size = DocxPt(10)

    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run("Expected student reasoning: ")
    run.bold = True
    run.font.size = DocxPt(10)
    p.add_run("\"Rainfall and Wind come from outside - we can't control weather. Dry Vegetation and Fire Spread happen inside the system as a result of the weather inputs.\"").font.size = DocxPt(10)

    doc.add_paragraph()

    add_section_header(doc, "Relationships", level=2)

    relationships = [
        ("Rainfall --> Dry Vegetation", "NEGATIVE (-)", "When it rains more, vegetation gets wet, so there's LESS dry vegetation. One goes up, the other goes down."),
        ("Dry Vegetation --> Fire Spread", "POSITIVE (+)", "More dry vegetation = more fuel = fire spreads more. Both go up together."),
        ("Wind --> Fire Spread", "POSITIVE (+)", "More wind = fire spreads faster and farther. Both go up together.")
    ]

    for connection, sign, explanation in relationships:
        p = doc.add_paragraph()
        run = p.add_run(connection + " = ")
        run.bold = True
        run.font.size = DocxPt(11)
        run = p.add_run(sign)
        run.bold = True
        run.font.color.rgb = DocxRGB(0x22, 0x8B, 0x22) if "+" in sign else DocxRGB(0xDC, 0x14, 0x3C)
        run.font.size = DocxPt(11)

        p = doc.add_paragraph()
        run = p.add_run("Expected reasoning: ")
        run.italic = True
        run.font.size = DocxPt(10)
        p.add_run(f"\"{explanation}\"").font.size = DocxPt(10)

    doc.add_paragraph()

    add_section_header(doc, "Simulation Observations", level=2)

    scenarios_answers = [
        ("Scenario 1: DROUGHT (Rainfall OFF)",
         "Dry Vegetation graph line should go UP (stays high or increases). Fire Spread should increase once dry vegetation is available. Students might say: \"When there's no rain, the plants dry out and stay dry. This makes more fuel for fires.\""),
        ("Scenario 2: WINDY DAY (Wind ON)",
         "Fire Spread should increase. Students might say: \"Wind pushes the fire and makes it spread faster. The graph shows fire spread going up when wind is on.\""),
        ("Scenario 3: FIRE SEASON (Rainfall OFF + Wind ON)",
         "Both Dry Vegetation and Fire Spread should be high/increasing. This represents the worst-case scenario. Students might say: \"This is the most dangerous! No rain means lots of dry fuel, and wind makes fire spread really fast. This is what happens in California every fall.\"")
    ]

    for scenario, answer in scenarios_answers:
        p = doc.add_paragraph()
        run = p.add_run(scenario)
        run.bold = True
        run.font.size = DocxPt(11)
        run.font.color.rgb = DOCX_BRAND_BLUE

        p = doc.add_paragraph()
        p.add_run(answer).font.size = DocxPt(10)
        doc.add_paragraph()

    add_section_header(doc, "Reflection Question Exemplars", level=2)

    reflections_answers = [
        ("Why does California burn at the same time every year?",
         "Strong response: \"California burns in fall because of how Earth's systems interact. During summer, there's no rain (hydrosphere), so vegetation dries out (biosphere) and becomes fuel. In fall, Santa Ana winds (atmosphere) blow dry and fast, spreading any fire quickly. It's the same every year because California has a Mediterranean climate with this predictable pattern.\"\n\nDeveloping response: \"There's no rain in summer so plants dry out, then wind makes fires spread.\""),
        ("How do Earth's systems interact?",
         "Strong response: \"The hydrosphere (water cycle) controls how wet plants are. When there's no rain, the biosphere (plants) dries out. The atmosphere (air/wind) then spreads fire through the dried plants. The geosphere (land) affects how fast fire moves uphill.\"\n\nDeveloping response: \"Water, air, land, and living things all affect each other. Less rain means dry plants, and wind spreads fire.\"")
    ]

    for question, answers in reflections_answers:
        p = doc.add_paragraph()
        run = p.add_run("Q: " + question)
        run.bold = True
        run.font.size = DocxPt(11)

        p = doc.add_paragraph()
        p.add_run(answers).font.size = DocxPt(10)
        doc.add_paragraph()

    doc.add_page_break()

    # ========== 6. STEM CHALLENGE TEACHER GUIDANCE ==========
    add_section_header(doc, "6. STEM Challenge: Teacher Guidance")

    add_section_header(doc, "How to Introduce", level=2)
    p = doc.add_paragraph()
    p.add_run("Say: ").font.size = DocxPt(11)
    p.add_run("\"You discovered that DRY VEGETATION is the FUEL for fire spread. Real engineers and firefighters use this knowledge to create FIREBREAKS - areas where they clear out all vegetation so fire has nothing to burn. When fire reaches a firebreak, it stops because there's no fuel! Now YOU will become Firebreak Engineers and design protection for a California town.\"").font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "Key Scientific Concepts", level=2)
    concepts = [
        ("Fire Triangle:", "Fire needs fuel, heat, and oxygen. Remove any one and fire stops. Firebreaks remove FUEL."),
        ("Wind Direction:", "Santa Ana winds blow from east to west (from the desert toward the coast). Firebreaks should be placed between the fire source (usually east) and what you want to protect."),
        ("Terrain Effects:", "Fire spreads faster uphill because flames preheat fuel above. If the town is in a valley, fires can race downhill toward it from multiple directions."),
        ("Strategic Placement:", "Best firebreak locations: between likely ignition points and the town, perpendicular to prevailing wind direction, along natural features like roads or rivers.")
    ]

    for concept, explanation in concepts:
        p = doc.add_paragraph()
        run = p.add_run(concept + " ")
        run.bold = True
        run.font.size = DocxPt(11)
        p.add_run(explanation).font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "Evaluating Student Designs", level=2)

    eval_criteria = [
        ("Strong Design:", "3 firebreaks placed east of town (toward wind source), considers terrain, explains reasoning using model concepts"),
        ("Developing Design:", "Some firebreaks correctly placed, partial reasoning about wind or fuel"),
        ("Needs Support:", "Random placement, no clear reasoning connected to the model")
    ]

    for level, criteria in eval_criteria:
        p = doc.add_paragraph()
        run = p.add_run(level + " ")
        run.bold = True
        run.font.size = DocxPt(11)
        p.add_run(criteria).font.size = DocxPt(11)

    doc.add_page_break()

    # ========== 7. DIFFERENTIATION & EXTENSIONS ==========
    add_section_header(doc, "7. Differentiation & Extensions")

    add_section_header(doc, "For Struggling Learners", level=2)
    support_strategies = [
        "Provide component cards that can be physically sorted before building in ModelIt!",
        "Use sentence frames: 'When ___ increases, ___ increases/decreases because ___.'",
        "Pre-build partial model and have students complete the relationships",
        "Pair with a stronger partner who can explain their thinking aloud"
    ]
    for strategy in support_strategies:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(strategy).font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "For Advanced Learners", level=2)
    extension_ideas = [
        "Add 2-3 additional components (Temperature, Humidity, Terrain) and justify the relationships",
        "Compare California fire systems to Australian bushfires or Amazon rainforest fires",
        "Research and model a feedback loop (fire releases CO2 --> climate change --> more fire?)",
        "Calculate firebreak costs and optimize design under budget constraints"
    ]
    for idea in extension_ideas:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(idea).font.size = DocxPt(11)

    doc.add_paragraph()

    add_section_header(doc, "Cross-Curricular Connections", level=2)
    connections = [
        ("Math:", "Calculate percentages (moisture content), area of firebreaks, cost analysis"),
        ("ELA:", "Write explanatory text about fire prevention using evidence from model"),
        ("Social Studies:", "Research how indigenous peoples used controlled burns; urban planning considerations")
    ]
    for subject, connection in connections:
        p = doc.add_paragraph()
        run = p.add_run(subject + " ")
        run.bold = True
        run.font.size = DocxPt(11)
        p.add_run(connection).font.size = DocxPt(11)

    doc.add_page_break()

    # ========== 8. COMMON MISCONCEPTIONS ==========
    add_section_header(doc, "8. Common Misconceptions & How to Address Them")

    misconceptions = [
        {
            "misconception": "All fires are bad and should be prevented",
            "reality": "Some ecosystems depend on fire! California chaparral evolved with fire. Controlled burns can actually reduce wildfire risk by clearing fuel.",
            "response": "Acknowledge that uncontrolled wildfires are dangerous, but introduce the idea that fire is a natural part of some ecosystems. This can be an extension discussion."
        },
        {
            "misconception": "Fire spread and fire starting are the same thing",
            "reality": "Our model focuses on SPREAD, not ignition. Fires can start from lightning, human activity, or power lines - our model shows what happens AFTER ignition.",
            "response": "Clarify that we're modeling how fires spread once they start, not what causes them to start in the first place."
        },
        {
            "misconception": "Relationships are always obvious",
            "reality": "Students may initially struggle with directionality. 'Rainfall affects dry vegetation' vs 'dry vegetation affects rainfall' are very different!",
            "response": "Use the 'If X increases, what happens to Y?' test consistently. Have students explain their reasoning aloud before drawing arrows."
        },
        {
            "misconception": "More components = better model",
            "reality": "Models should be as simple as possible while still being useful. Too many components can make the model confusing without adding insight.",
            "response": "Praise simplicity! A clear 4-component model that students understand is better than a complex 10-component model they can't explain."
        }
    ]

    for item in misconceptions:
        p = doc.add_paragraph()
        run = p.add_run("Misconception: ")
        run.bold = True
        run.font.size = DocxPt(11)
        run.font.color.rgb = DocxRGB(0xDC, 0x14, 0x3C)
        p.add_run(item['misconception']).font.size = DocxPt(11)

        p = doc.add_paragraph()
        run = p.add_run("Reality: ")
        run.bold = True
        run.font.size = DocxPt(10)
        p.add_run(item['reality']).font.size = DocxPt(10)

        p = doc.add_paragraph()
        run = p.add_run("How to respond: ")
        run.italic = True
        run.font.size = DocxPt(10)
        p.add_run(item['response']).font.size = DocxPt(10)

        doc.add_paragraph()

    doc.save(output_path)
    print(f"[OK] Created: {output_path}")
    return output_path


# ============================================
# MAIN
# ============================================

if __name__ == "__main__":
    ensure_output_dir()

    print("\n" + "=" * 60)
    print("ModelIt! Lesson Materials Generator v4")
    print(f"{LESSON_ID}: {LESSON_TITLE}")
    print("With AI Image Generation & Aspect Ratio Preservation")
    print("=" * 60)

    # Check if images already exist
    existing_images = {}
    image_files = {
        'cover': 'cover-students-collaborating.png',
        'landscape': 'california-fire-season.png',
        'discussion': 'students-discussion.png',
        'modeling': 'student-modeling.png',
        'stem': 'stem-challenge.png'
    }

    all_exist = True
    for key, filename in image_files.items():
        path = os.path.join(IMAGES_DIR, filename)
        if os.path.exists(path):
            existing_images[key] = path
        else:
            all_exist = False

    if all_exist:
        print("\n[OK] Using existing images (all 5 found)")
        images = existing_images
    else:
        print("\n[!] Some images missing - generating new ones")
        images = generate_all_lesson_images()

    print("\n" + "=" * 60)
    print("CREATING MATERIALS")
    print("=" * 60 + "\n")

    # Create materials
    ppt_path = create_powerpoint(images)
    activity_path = create_student_activity_pack()
    guide_path = create_teachers_guide()

    print("\n" + "=" * 60)
    print("COMPLETE!")
    print("=" * 60)
    print(f"\nLocation: {OUTPUT_DIR}")
    print(f"\nFiles created:")
    print(f"  - {LESSON_ID}-Student-Presentation.pptx")
    print(f"  - {LESSON_ID}-Student-Activity-Pack.docx")
    print(f"  - {LESSON_ID}-Teachers-Guide.docx")
    print(f"\nImages used: {len([v for v in images.values() if v])}")
    if total_image_cost > 0:
        print(f"Total image cost: ${total_image_cost:.4f}")
    print("=" * 60)
