"""
ModelIt Branded Lesson Materials Generator
Creates PowerPoint, Student Activity Pack, and Teacher's Guide
matching the official TPT branding style.

Brand Elements:
- Blue diagonal geometric corners (navy, medium blue, light blue layers)
- ModelIt! logo with molecule/network icon
- Science doodle border for B&W worksheet pages
- Student Name/Date on every page
- STEM Challenge included
- Recording Page included
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, black, white, Color
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT
import os

# ============================================
# BRAND COLORS
# ============================================
# From TPT samples - layered blue diagonal design
NAVY = RGBColor(0x0D, 0x1B, 0x2A)        # Darkest blue
BRAND_BLUE = RGBColor(0x1A, 0x47, 0x80)   # Deep blue
MID_BLUE = RGBColor(0x2E, 0x86, 0xAB)     # Medium blue
LIGHT_BLUE = RGBColor(0x7E, 0xC8, 0xE3)   # Light accent
SKY_BLUE = RGBColor(0x5D, 0xB7, 0xDE)     # Logo background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)

# PDF Colors
PDF_NAVY = HexColor('#0D1B2A')
PDF_BRAND_BLUE = HexColor('#1A4780')
PDF_MID_BLUE = HexColor('#2E86AB')
PDF_LIGHT_BLUE = HexColor('#7EC8E3')
PDF_SKY_BLUE = HexColor('#5DB7DE')

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'materials', 'G05-L01')

def ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(os.path.join(OUTPUT_DIR, 'images'), exist_ok=True)

# ============================================
# POWERPOINT - BRANDED DESIGN
# ============================================

def add_diagonal_corners(slide, prs, position='both'):
    """Add the signature blue diagonal corner design"""
    w = prs.slide_width
    h = prs.slide_height

    if position in ['top', 'both']:
        # Top-left diagonal stripes
        # Navy (darkest, largest)
        shape1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(-0.5), Inches(-0.5), Inches(4), Inches(1.5))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = NAVY
        shape1.line.fill.background()
        shape1.rotation = 0

        # Medium blue (middle)
        shape2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(2.5), Inches(0), Inches(3), Inches(0.8))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = MID_BLUE
        shape2.line.fill.background()

        # Light blue (accent stripe)
        shape3 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(3.5), Inches(0.5), Inches(2.5), Inches(0.4))
        shape3.fill.solid()
        shape3.fill.fore_color.rgb = LIGHT_BLUE
        shape3.line.fill.background()

    if position in ['bottom', 'both']:
        # Bottom-right diagonal stripes
        # Navy (darkest)
        shape4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(6.5), Inches(6), Inches(4), Inches(1.5))
        shape4.fill.solid()
        shape4.fill.fore_color.rgb = NAVY
        shape4.line.fill.background()
        shape4.rotation = 180

        # Medium blue
        shape5 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(4.5), Inches(6.7), Inches(3), Inches(0.8))
        shape5.fill.solid()
        shape5.fill.fore_color.rgb = MID_BLUE
        shape5.line.fill.background()

        # Light blue
        shape6 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
            Inches(4), Inches(6.3), Inches(2.5), Inches(0.4))
        shape6.fill.solid()
        shape6.fill.fore_color.rgb = LIGHT_BLUE
        shape6.line.fill.background()

def add_modelit_logo(slide, x, y, width=3):
    """Add ModelIt! logo box"""
    # Blue background rectangle
    logo_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(0.9))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = SKY_BLUE
    logo_bg.line.fill.background()

    # Molecule icon placeholder (small white circle cluster)
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x + 0.15), Inches(y + 0.25), Inches(0.4), Inches(0.4))
    icon.fill.solid()
    icon.fill.fore_color.rgb = WHITE
    icon.line.fill.background()

    # ModelIt! text
    text_box = slide.shapes.add_textbox(
        Inches(x + 0.6), Inches(y + 0.15), Inches(width - 0.8), Inches(0.6))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ModelIt!"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Comic Sans MS"  # Handwritten style

def add_side_dots(slide, y_start, count=7):
    """Add decorative dots on sides like TPT design"""
    # Left side dots
    for i in range(count):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.3), Inches(y_start + i * 0.25), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = NAVY
        dot.line.fill.background()

    # Right side dots
    for i in range(count):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(9.6), Inches(y_start + i * 0.25), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = NAVY
        dot.line.fill.background()

def add_image_placeholder(slide, x, y, width, height, description):
    """Add a placeholder box for images with description"""
    # Placeholder box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFA)
    box.line.color.rgb = MID_BLUE
    box.line.width = Pt(2)

    # Description text
    text_box = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(y + height/2 - 0.3),
        Inches(width - 0.4), Inches(0.6))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[IMAGE: {description}]"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

def create_cover_slide(prs, title, subtitle):
    """Create branded cover slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add diagonal corners
    add_diagonal_corners(slide, prs)

    # Add ModelIt logo
    add_modelit_logo(slide, 3.5, 1.8)

    # Add side dots
    add_side_dots(slide, 3.5)

    # Title section header
    header = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Student PowerPoint"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Lesson title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Image placeholder for cover
    add_image_placeholder(slide, 2.5, 4.8, 5, 1.8,
        "Diverse 5th graders looking at California golden hills with distant wildfire smoke. 70% students of color (Black, Latino, Asian). Age-appropriate, educational.")

    return slide

def create_content_slide(prs, title, bullets, image_desc=None):
    """Create a content slide with optional image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Top diagonal only
    add_diagonal_corners(slide, prs, 'top')

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Content area
    if image_desc:
        # Split layout - text left, image right
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4.5))
        add_image_placeholder(slide, 5.2, 2, 4.3, 3.5, image_desc)
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(4), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "ModelIt! | G05-L01"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    return slide

def create_activity_slide(prs, activity_num, activity_name, steps, platform_placeholder=True):
    """Create an activity instruction slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_diagonal_corners(slide, prs, 'top')

    # Activity badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.4), Inches(1.1), Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = MID_BLUE
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(Inches(0.4), Inches(1.25), Inches(0.7), Inches(0.4))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = str(activity_num)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Activity name
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(1.15), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = activity_name
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Steps
    if platform_placeholder:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5), Inches(4.5))
        add_image_placeholder(slide, 5.7, 2, 4, 3,
            "PLATFORM SCREENSHOT - Team will add actual ModelIt interface screenshot")
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, step in enumerate(steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

    return slide

def create_question_slide(prs, question, subtext=""):
    """Create a driving question slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MID_BLUE
    bg.line.fill.background()

    # Large question mark
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1.5), Inches(2))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "?"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    # Question text
    q_text = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = q_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtext:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_powerpoint():
    """Create the full branded student PowerPoint"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    create_cover_slide(prs,
        "When Trees Become Matches",
        "California's Burning Season and the Earth Systems That Fuel It")

    # Slide 2: Hook Image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')
    add_image_placeholder(slide, 0.5, 1.2, 9, 5.5,
        "Split image: Beautiful green California forest on left / Same forest with orange fire-smoke sky on right. Dramatic contrast showing before/after of wildfire impact.")
    caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf = caption.text_frame
    p = tf.paragraphs[0]
    p.text = "Same state. Same trees. Every single year. Why?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Slide 3: Driving Question
    create_question_slide(prs,
        "How do Earth's systems work together to create perfect fire conditions?",
        "Atmosphere + Biosphere + Hydrosphere = 🔥")

    # Slide 4: What We're Building
    create_content_slide(prs, "What We're Building Today", [
        "A cause-and-effect MODEL of wildfire conditions",
        "Discover how Earth's systems connect",
        "Predict what happens when conditions change",
        "Think like a Wildfire Data Scientist"
    ], "Diverse students (Black girl, Latino boy, Asian girl) collaborating on computers in science classroom. Engaged, curious expressions. Age 10-11.")

    # Slide 5: Components
    create_content_slide(prs, "The Parts of Our Fire System", [
        "EXTERNAL — things we CAN'T control",
        "   • Rainfall (weather decides)",
        "   • Wind (nature decides)",
        "",
        "INTERNAL — things that CHANGE",
        "   • Dry Vegetation (changes based on rain)",
        "   • Fire Spread (what we're trying to understand)"
    ])

    # Slide 6: Activity 1
    create_activity_slide(prs, 1, "LOCATE — Build Your Fire System", [
        "STEP 1: Look at the components in ModelIt",
        "",
        "STEP 2: SORT them into EXTERNAL and INTERNAL",
        "   • Your teacher will show you how!",
        "",
        "STEP 3: Click the PLUS (+) button to add each",
        "",
        "STEP 4: You should have 4 components on your canvas"
    ])

    # Slide 7: Activity 2
    create_activity_slide(prs, 2, "ESTABLISH — Connect the Relationships", [
        "STEP 1: Click the CONNECT icon (top left)",
        "",
        "STEP 2: Draw arrows between components:",
        "   • Rainfall → Dry Vegetation",
        "   • Dry Vegetation → Fire Spread",
        "   • Wind → Fire Spread",
        "",
        "STEP 3: Set each as + (positive) or − (negative)",
        "",
        "Ask: When THIS goes up, does THAT go up or down?"
    ])

    # Slide 8: Relationships Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Chain Reaction"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    add_image_placeholder(slide, 1, 2, 8, 3.5,
        "DIAGRAM: Rainfall —(−)→ Dry Vegetation —(+)→ Fire Spread ←(+)— Wind. Show arrows with +/- labels. Clean, educational style matching ModelIt aesthetic.")

    caption = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    tf = caption.text_frame
    p = tf.paragraphs[0]
    p.text = "When rain STOPS, the dominoes fall toward FIRE"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = MID_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Slide 9: Activity 3
    create_activity_slide(prs, 3, "VISUALIZE — Run Your Fire Model", [
        "STEP 1: Click the PLAY button (top left)",
        "",
        "STEP 2: Watch the graph — see components change",
        "",
        "STEP 3: CREATE A DROUGHT:",
        "   • Find the LOCK icon next to Rainfall",
        "   • Lock Rainfall to OFF (0%)",
        "",
        "STEP 4: Watch Dry Vegetation and Fire Spread rise!",
        "",
        "STEP 5: Lock WIND to ON (100%) — Santa Ana winds!"
    ])

    # Slide 10: Graph Reading
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Reading Your Simulation"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    add_image_placeholder(slide, 0.5, 2, 9, 4,
        "PLATFORM SCREENSHOT: Graph showing Rainfall at 0%, Dry Vegetation climbing to high %, Fire Spread climbing. Show colorful lines on graph panel.")

    caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = caption.text_frame
    p = tf.paragraphs[0]
    p.text = "You just simulated California's fire season!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MID_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Slide 11: Activity 4
    create_activity_slide(prs, 4, "EXTEND — Play, Research, Expand", [
        "PLAY TIME! Experiment with your model:",
        "   • What if rain comes back?",
        "   • What creates the WORST fire conditions?",
        "",
        "RESEARCH: What's MISSING from your model?",
        "   • Temperature / Heat waves?",
        "   • Humans (campfires, power lines)?",
        "   • Firefighters?",
        "",
        "ADD one new component based on your research!"
    ], platform_placeholder=False)

    # Slide 12: Earth Systems
    create_content_slide(prs, "Earth Systems Working Together", [
        "ATMOSPHERE — weather, wind, temperature",
        "",
        "BIOSPHERE — plants, trees, vegetation",
        "",
        "HYDROSPHERE — rain, drought, water cycle",
        "",
        "Wildfires happen where all three OVERLAP",
        "Your model shows how they CONNECT"
    ], "Venn diagram showing Atmosphere, Biosphere, Hydrosphere overlapping. Fire icon in center where all three meet. Blue/green color scheme.")

    # Slide 13: Career Connection
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_diagonal_corners(slide, prs, 'top')

    # Orange accent for career
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1), Inches(0.15), Inches(6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0xE6, 0x55, 0x00)
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "🔥 Career Connection: Wildfire Data Scientist"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    content = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5), Inches(4))
    tf = content.text_frame
    tf.word_wrap = True

    bullets = [
        "Builds computer models just like yours — but BIGGER",
        "Predicts where fires will start BEFORE they happen",
        "Analyzes weather, vegetation, wind, and climate data",
        "Helps position firefighters and warn communities",
        "",
        "💰 SALARY: $95,000 - $150,000+ per year",
        "",
        "Your fire model? That's step one toward this career."
    ]

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}" if bullet and not bullet.startswith("💰") and not bullet.startswith("Your") else bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    add_image_placeholder(slide, 5.7, 2, 4, 3.5,
        "Black or Latino data scientist (woman, 30s) at computer with fire prediction map on screen. Professional setting, confident expression. Inspirational role model image.")

    # Slide 14: Reflection
    create_question_slide(prs,
        "What story does YOUR model tell?",
        "Think-Pair-Share: Explain your model to a partner")

    # Save
    output_path = os.path.join(OUTPUT_DIR, 'G05-L01-Student-Presentation-BRANDED.pptx')
    prs.save(output_path)
    print(f"[OK] Created: {output_path}")
    return output_path


# ============================================
# STUDENT ACTIVITY PACK - B&W WITH BORDERS
# ============================================

class BrandedWorksheetPDF:
    """Creates B&W worksheets with science doodle border styling"""

    def __init__(self, filename, lesson_id="G05-L01"):
        self.filename = filename
        self.lesson_id = lesson_id
        self.c = canvas.Canvas(filename, pagesize=letter)
        self.width, self.height = letter
        self.margin = 0.75 * inch
        self.y = self.height - self.margin - 0.5 * inch
        self.page_num = 1
        self.content_start_y = self.height - 1.2 * inch

    def draw_doodle_border(self):
        """Draw science doodle border (simplified version - rectangles as placeholders)"""
        # For actual production, these would be replaced with actual doodle images
        # Drawing a decorative border frame
        self.c.setStrokeColor(black)
        self.c.setLineWidth(1.5)

        # Outer border
        self.c.rect(0.4 * inch, 0.4 * inch,
                   self.width - 0.8 * inch, self.height - 0.8 * inch)

        # Inner content area
        self.c.setLineWidth(0.5)
        self.c.rect(0.6 * inch, 0.6 * inch,
                   self.width - 1.2 * inch, self.height - 1.2 * inch)

        # Corner decorations (simple shapes to represent doodles)
        self.c.setFillColor(black)

        # Small circles in corners to suggest doodle style
        corners = [
            (0.5 * inch, self.height - 0.5 * inch),
            (self.width - 0.5 * inch, self.height - 0.5 * inch),
            (0.5 * inch, 0.5 * inch),
            (self.width - 0.5 * inch, 0.5 * inch)
        ]
        for x, y in corners:
            self.c.circle(x, y, 3, fill=1)

    def add_header(self, title):
        """Add page header with name/date line"""
        self.draw_doodle_border()

        # Title in stylized font
        self.c.setFont("Helvetica-Bold", 16)
        self.c.drawCentredString(self.width / 2, self.height - 0.9 * inch, title)

        # Name and Date line
        self.c.setFont("Helvetica", 10)
        self.c.drawString(0.75 * inch, self.height - 1.15 * inch,
                         "Student Name: _________________________________")
        self.c.drawString(5.5 * inch, self.height - 1.15 * inch,
                         "Date: _______________")

        self.y = self.content_start_y - 0.2 * inch

    def add_footer(self):
        """Add page footer"""
        self.c.setFont("Helvetica", 8)
        self.c.setFillColor(Color(0.4, 0.4, 0.4))
        self.c.drawString(0.75 * inch, 0.7 * inch, f"{self.lesson_id} | Page {self.page_num}")
        self.c.drawRightString(self.width - 0.75 * inch, 0.7 * inch,
                              "ModelIt! | Systems Thinking for K-12")
        self.c.setFillColor(black)

    def add_section_title(self, title):
        """Add a section title"""
        if self.y < 2 * inch:
            self.new_page()
        self.c.setFont("Helvetica-Bold", 14)
        self.c.drawString(self.margin, self.y, title)
        self.y -= 8
        self.c.setLineWidth(1)
        self.c.line(self.margin, self.y, self.margin + 250, self.y)
        self.y -= 20

    def add_text(self, text, bold=False, indent=0, size=11):
        """Add a line of text"""
        if self.y < 1.2 * inch:
            self.new_page()
        font = "Helvetica-Bold" if bold else "Helvetica"
        self.c.setFont(font, size)
        self.c.drawString(self.margin + indent, self.y, text)
        self.y -= 16

    def add_line(self, label="", width=400):
        """Add a fill-in line"""
        if self.y < 1.2 * inch:
            self.new_page()
        self.c.setFont("Helvetica", 11)
        if label:
            self.c.drawString(self.margin, self.y, label)
            line_start = self.margin + self.c.stringWidth(label, "Helvetica", 11) + 5
        else:
            line_start = self.margin
        self.c.line(line_start, self.y - 2, self.margin + width, self.y - 2)
        self.y -= 24

    def add_lines(self, count=3, width=450):
        """Add multiple blank lines"""
        for _ in range(count):
            if self.y < 1.2 * inch:
                self.new_page()
            self.c.line(self.margin, self.y, self.margin + width, self.y)
            self.y -= 24

    def add_box(self, height=100, label=""):
        """Add a drawing/answer box"""
        if self.y - height < 1.2 * inch:
            self.new_page()
        self.c.setStrokeColor(black)
        self.c.setLineWidth(1)
        self.c.rect(self.margin, self.y - height, self.width - 2 * self.margin, height)
        if label:
            self.c.setFont("Helvetica-Oblique", 9)
            self.c.setFillColor(Color(0.5, 0.5, 0.5))
            self.c.drawString(self.margin + 5, self.y - 14, label)
            self.c.setFillColor(black)
        self.y -= height + 18

    def add_checkbox_item(self, text):
        """Add a checkbox item"""
        if self.y < 1.2 * inch:
            self.new_page()
        self.c.rect(self.margin, self.y - 10, 12, 12)
        self.c.setFont("Helvetica", 11)
        self.c.drawString(self.margin + 18, self.y - 8, text)
        self.y -= 22

    def add_space(self, points=15):
        """Add vertical space"""
        self.y -= points

    def new_page(self, title=""):
        """Start a new page"""
        self.add_footer()
        self.c.showPage()
        self.page_num += 1
        self.y = self.height - self.margin
        if title:
            self.add_header(title)
        else:
            self.draw_doodle_border()
            self.y = self.content_start_y

    def save(self):
        """Save the PDF"""
        self.add_footer()
        self.c.save()
        print(f"[OK] Created: {self.filename}")


def create_activity_pack():
    """Create the branded student activity pack"""
    output_path = os.path.join(OUTPUT_DIR, 'G05-L01-Activity-Pack-BRANDED.pdf')
    pdf = BrandedWorksheetPDF(output_path)

    # ========== PAGE 1: COVER ==========
    pdf.c.setFont("Helvetica-Bold", 28)
    pdf.c.drawCentredString(pdf.width / 2, pdf.height - 2 * inch, "ModelIt!")
    pdf.c.setFont("Helvetica-Bold", 20)
    pdf.c.drawCentredString(pdf.width / 2, pdf.height - 2.6 * inch, "Printable Activity Pack")
    pdf.c.setFont("Helvetica", 16)
    pdf.c.drawCentredString(pdf.width / 2, pdf.height - 3.2 * inch,
                           "When Trees Become Matches")
    pdf.c.setFont("Helvetica-Oblique", 12)
    pdf.c.drawCentredString(pdf.width / 2, pdf.height - 3.6 * inch,
                           "California's Burning Season and the Earth Systems That Fuel It")

    # Cover image placeholder
    pdf.c.setStrokeColor(black)
    pdf.c.rect(2 * inch, 3 * inch, 4.5 * inch, 3 * inch)
    pdf.c.setFont("Helvetica-Oblique", 10)
    pdf.c.drawCentredString(pdf.width / 2, 4.3 * inch,
                           "[IMAGE: California hills with wildfire smoke]")

    pdf.add_footer()
    pdf.c.showPage()
    pdf.page_num += 1

    # ========== PAGE 2: TABLE OF CONTENTS ==========
    pdf.add_header("TABLE OF CONTENTS")
    pdf.add_space(20)
    pdf.add_text("1. Scoring Rubric", size=14)
    pdf.add_text("2. Anecdotal Notes Sheet", size=14)
    pdf.add_text("3. Model Recording Page", size=14)
    pdf.add_text("4. Simulation Observations", size=14)
    pdf.add_text("5. Research & Extend", size=14)
    pdf.add_text("6. STEM Challenge: Firebreak Engineers", size=14)
    pdf.add_text("7. STEM Challenge: Teacher Notes", size=14)
    pdf.new_page()

    # ========== PAGE 3: SCORING RUBRIC ==========
    pdf.add_header("SCORING RUBRIC")
    pdf.add_space(10)

    # Create rubric table manually
    pdf.c.setFont("Helvetica-Bold", 10)
    col_widths = [1.2 * inch, 1.5 * inch, 1.5 * inch, 1.5 * inch, 1.3 * inch]
    headers = ["Criteria", "4-Extending", "3-Proficient", "2-Developing", "1-Beginning"]

    x = pdf.margin
    for i, header in enumerate(headers):
        pdf.c.drawString(x + 5, pdf.y, header)
        x += col_widths[i]

    pdf.y -= 15
    pdf.c.setLineWidth(1)
    pdf.c.line(pdf.margin, pdf.y, pdf.width - pdf.margin, pdf.y)
    pdf.y -= 15

    # Rubric rows
    rubric_data = [
        ["Model\nConstruction", "Builds model +\nextra components", "Builds and labels\nall parts", "Most parts there,\nsome missing", "Model incomplete\nor incorrect"],
        ["Connections\n(Arrows)", "All arrows +\nextra correct", "All required\narrows correct", "Some arrows\nmissing/incorrect", "Few or no\ncorrect arrows"],
        ["Simulation\nUse", "Runs several tests,\nexplains results", "Runs simulation,\nbasic explanation", "Runs with help,\nunclear explanation", "Cannot run or\nexplain"],
        ["Participation", "On task, asks\nquestions, helps", "On task,\nfollows directions", "Sometimes off task,\nneeds reminders", "Rarely participates"]
    ]

    pdf.c.setFont("Helvetica", 8)
    for row in rubric_data:
        x = pdf.margin
        max_lines = max(len(cell.split('\n')) for cell in row)
        for i, cell in enumerate(row):
            lines = cell.split('\n')
            for j, line in enumerate(lines):
                pdf.c.drawString(x + 3, pdf.y - j * 10, line)
            x += col_widths[i]
        pdf.y -= max_lines * 12 + 8
        pdf.c.line(pdf.margin, pdf.y + 5, pdf.width - pdf.margin, pdf.y + 5)

    pdf.add_space(20)
    pdf.add_text("Total Score: ___________", bold=True)
    pdf.add_text("14-16 = Extending | 10-13 = Proficient | 7-9 = Developing | 4-6 = Beginning", size=9)

    pdf.new_page()

    # ========== PAGE 4: ANECDOTAL NOTES ==========
    pdf.add_header("ANECDOTAL NOTES SHEET")
    pdf.add_space(10)

    pdf.add_text("Unit: When Trees Become Matches (Earth Systems)", bold=True)
    pdf.add_space(15)

    pdf.add_text("Verbal Reasoning:", bold=True)
    pdf.add_lines(2)

    pdf.add_text("Concept Cognition:", bold=True)
    pdf.add_lines(2)

    pdf.add_text("Peer Collaboration:", bold=True)
    pdf.add_lines(2)

    pdf.add_space(10)
    pdf.add_text("Checklist:", bold=True)
    pdf.add_checkbox_item("Used Graph Panel correctly")
    pdf.add_checkbox_item("Ran simulation with different scenarios")
    pdf.add_checkbox_item("Identified cause-effect flow")
    pdf.add_checkbox_item("Asked questions or made real-world links")
    pdf.add_checkbox_item("Added research-based component")

    pdf.add_space(15)
    pdf.add_text("Additional Notes:", bold=True)
    pdf.add_box(100)

    pdf.new_page()

    # ========== PAGE 5: MODEL RECORDING PAGE ==========
    pdf.add_header("MODEL RECORDING PAGE")
    pdf.add_space(10)

    pdf.add_text("Draw your completed model below. Include all components and arrows.", bold=True)
    pdf.add_text("Label each connection as + (positive) or − (negative).")
    pdf.add_space(10)

    pdf.add_box(200, "Draw your fire system model here")

    pdf.add_text("List your components:", bold=True)
    pdf.add_text("EXTERNAL (can't control):", indent=20)
    pdf.add_line("1. ", 200)
    pdf.add_line("2. ", 200)

    pdf.add_text("INTERNAL (changes based on other things):", indent=20)
    pdf.add_line("1. ", 200)
    pdf.add_line("2. ", 200)

    pdf.new_page()

    # ========== PAGE 6: SIMULATION OBSERVATIONS ==========
    pdf.add_header("SIMULATION OBSERVATIONS")
    pdf.add_space(10)

    pdf.add_text("BASELINE TEST (Everything ON)", bold=True)
    pdf.add_text("What happens when the simulation runs normally?")
    pdf.add_lines(2)

    pdf.add_space(10)
    pdf.add_text("DROUGHT TEST (Rainfall OFF)", bold=True)
    pdf.add_text("When I turned OFF Rainfall, this happened:")
    pdf.add_text("• Dry Vegetation:    went UP  /  went DOWN  /  stayed same", indent=20)
    pdf.add_text("• Fire Spread:       went UP  /  went DOWN  /  stayed same", indent=20)

    pdf.add_space(10)
    pdf.add_text("WIND TEST (Wind ON)", bold=True)
    pdf.add_text("When I turned ON Wind (Santa Ana winds), this happened:")
    pdf.add_text("• Fire Spread:       went UP  /  went DOWN  /  stayed same", indent=20)

    pdf.add_space(10)
    pdf.add_text("WORST CONDITIONS", bold=True)
    pdf.add_text("The WORST fire conditions in my model are when:")
    pdf.add_lines(2)

    pdf.add_space(10)
    pdf.add_text("Sketch what your graph looked like:", bold=True)
    pdf.add_box(120, "Sketch the graph lines here")

    pdf.new_page()

    # ========== PAGE 7: RESEARCH & EXTEND ==========
    pdf.add_header("RESEARCH & EXTEND")
    pdf.add_space(10)

    pdf.add_text("What's MISSING from your model? Research one new component!", bold=True)
    pdf.add_space(10)

    pdf.add_line("NEW COMPONENT I want to add: ", 300)
    pdf.add_space(5)
    pdf.add_text("Is it EXTERNAL or INTERNAL? (circle one)")
    pdf.add_space(5)
    pdf.add_line("What does it connect to? ", 300)
    pdf.add_space(5)
    pdf.add_text("Is the relationship POSITIVE (+) or NEGATIVE (−)? _____")
    pdf.add_space(5)
    pdf.add_text("Why? Explain your reasoning:")
    pdf.add_lines(3)

    pdf.add_space(10)
    pdf.add_text("After adding it, my simulation showed:")
    pdf.add_lines(3)

    pdf.add_space(15)
    pdf.add_section_title("REFLECTION")
    pdf.add_text("Explain why California has wildfires every year (use your model!):")
    pdf.add_lines(4)

    pdf.add_text("Which Earth systems interact to create fire conditions? (circle all)")
    pdf.add_text("ATMOSPHERE      BIOSPHERE      HYDROSPHERE      GEOSPHERE", indent=40)

    pdf.new_page()

    # ========== PAGE 8: STEM CHALLENGE ==========
    pdf.add_header("STEM CHALLENGE: FIREBREAK ENGINEERS")
    pdf.add_space(10)

    pdf.add_text("THE CHALLENGE", bold=True, size=14)
    pdf.add_space(5)
    pdf.add_text("Firefighters create FIREBREAKS — gaps in vegetation that stop fire spread.")
    pdf.add_text("Your mission: Design a firebreak plan for a California community!")

    pdf.add_space(15)
    pdf.add_text("SCENARIO:", bold=True)
    pdf.add_text("A small town is surrounded by dry brush and trees. Based on your model,")
    pdf.add_text("you know that DRY VEGETATION + WIND = major fire spread. Where would")
    pdf.add_text("you put firebreaks to protect the town?")

    pdf.add_space(15)
    pdf.add_text("YOUR FIREBREAK PLAN:", bold=True)
    pdf.add_box(150, "Draw a simple map showing: Town (square) + Trees (circles) + Your firebreaks (lines)")

    pdf.add_text("Explain your design:", bold=True)
    pdf.add_text("Why did you put firebreaks in those locations?")
    pdf.add_lines(3)

    pdf.add_text("Connection to your model:", bold=True)
    pdf.add_text("How does your firebreak affect the FIRE SPREAD component?")
    pdf.add_lines(2)

    pdf.new_page()

    # ========== PAGE 9: STEM CHALLENGE TEACHER NOTES ==========
    pdf.add_header("STEM CHALLENGE: TEACHER NOTES")
    pdf.add_space(10)

    pdf.add_text("FIREBREAK ENGINEERS", bold=True, size=14)
    pdf.add_space(10)

    pdf.add_text("Big Idea:", bold=True)
    pdf.add_text("Firebreaks work by removing FUEL from the fire system. Students apply")
    pdf.add_text("their model understanding to a real-world engineering solution.")

    pdf.add_space(10)
    pdf.add_text("Materials Needed:", bold=True)
    pdf.add_text("• Paper and pencil (already provided in activity pack)")
    pdf.add_text("• Optional: colored pencils for map drawing")
    pdf.add_text("• NO PREP REQUIRED!")

    pdf.add_space(10)
    pdf.add_text("Implementation:", bold=True)
    pdf.add_text("1. Review the model: Dry Vegetation → Fire Spread")
    pdf.add_text("2. Introduce firebreaks as 'removing the fuel'")
    pdf.add_text("3. Students design individually or in pairs (10-15 min)")
    pdf.add_text("4. Share designs — discuss trade-offs")

    pdf.add_space(10)
    pdf.add_text("Discussion Questions:", bold=True)
    pdf.add_text("• Why can't we just remove ALL vegetation?")
    pdf.add_text("• Where do real firefighters put firebreaks?")
    pdf.add_text("• How does WIND affect where firebreaks should go?")

    pdf.add_space(10)
    pdf.add_text("Connection Made:", bold=True)
    pdf.add_text("\"If we reduce DRY VEGETATION in key areas, we BREAK the chain")
    pdf.add_text("that leads to FIRE SPREAD — just like our model showed!\"")

    pdf.add_space(10)
    pdf.add_text("Optional Extension:", bold=True)
    pdf.add_text("Add 'Firebreaks' as a new component in ModelIt with a NEGATIVE")
    pdf.add_text("relationship to Fire Spread. Run the simulation to see the effect!")

    pdf.save()
    return output_path


# ============================================
# TEACHER'S GUIDE - BRANDED
# ============================================

def create_teachers_guide():
    """Create branded teacher's guide (color PDF)"""
    output_path = os.path.join(OUTPUT_DIR, 'G05-L01-Teachers-Guide-BRANDED.pdf')

    styles = getSampleStyleSheet()

    # Custom styles matching brand
    title_style = ParagraphStyle(
        'BrandTitle',
        parent=styles['Title'],
        fontSize=24,
        textColor=PDF_NAVY,
        spaceAfter=20
    )

    heading_style = ParagraphStyle(
        'BrandHeading',
        parent=styles['Heading1'],
        fontSize=16,
        textColor=PDF_BRAND_BLUE,
        spaceBefore=15,
        spaceAfter=8
    )

    body_style = ParagraphStyle(
        'BrandBody',
        parent=styles['Normal'],
        fontSize=11,
        spaceAfter=8,
        leading=14
    )

    bullet_style = ParagraphStyle(
        'BrandBullet',
        parent=styles['Normal'],
        fontSize=11,
        leftIndent=20,
        spaceAfter=4
    )

    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        rightMargin=0.75*inch,
        leftMargin=0.75*inch,
        topMargin=0.75*inch,
        bottomMargin=0.75*inch
    )

    story = []

    # Title
    story.append(Paragraph("ModelIt! Teacher Guide", title_style))
    story.append(Paragraph("G05-L01: When Trees Become Matches", heading_style))
    story.append(Paragraph("<i>California's Burning Season and the Earth Systems That Fuel It</i>", body_style))
    story.append(Spacer(1, 20))

    # Overview box
    overview_data = [
        ['Grade Level', '5th Grade'],
        ['NGSS Standard', '5-ESS2-1'],
        ['Time', '40-45 minutes (can split across 2 days)'],
        ['Materials', 'Student devices with ModelIt access'],
        ['Prep Required', 'None — lesson is self-contained'],
    ]
    overview_table = Table(overview_data, colWidths=[1.8*inch, 4.2*inch])
    overview_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), PDF_LIGHT_BLUE),
        ('TEXTCOLOR', (0, 0), (0, -1), PDF_NAVY),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTNAME', (1, 0), (1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
        ('ALIGN', (1, 0), (1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('GRID', (0, 0), (-1, -1), 0.5, PDF_MID_BLUE),
        ('PADDING', (0, 0), (-1, -1), 8),
    ]))
    story.append(overview_table)
    story.append(Spacer(1, 20))

    # Learning Objectives
    story.append(Paragraph("LEARNING OBJECTIVES", heading_style))
    objectives = [
        "Model how lack of rainfall affects vegetation",
        "Trace cause-and-effect relationships between drought and fire conditions",
        "Explain how Earth systems (atmosphere, biosphere, hydrosphere) interact",
        "Predict what happens when one part of the system changes"
    ]
    for obj in objectives:
        story.append(Paragraph(f"• {obj}", bullet_style))
    story.append(Spacer(1, 15))

    # LEVER Alignment
    story.append(Paragraph("LEVER FRAMEWORK ALIGNMENT", heading_style))
    lever_data = [
        ['Activity', 'LEVER Phase', 'Time', 'What Students Do'],
        ['Activity 1', 'LOCATE', '8-10 min', 'Sort and add components'],
        ['Activity 2', 'ESTABLISH', '8-10 min', 'Connect with arrows (+/−)'],
        ['Activity 3', 'VISUALIZE & EVALUATE', '10-12 min', 'Run simulations, analyze graphs'],
        ['Activity 4', 'REVISE & EXTEND', '10-15 min', 'Research and expand model'],
    ]
    lever_table = Table(lever_data, colWidths=[1.2*inch, 1.5*inch, 1*inch, 2.8*inch])
    lever_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), PDF_BRAND_BLUE),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('GRID', (0, 0), (-1, -1), 0.5, PDF_MID_BLUE),
        ('PADDING', (0, 0), (-1, -1), 6),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#F0F5FA')]),
    ]))
    story.append(lever_table)
    story.append(Spacer(1, 15))

    # Answer Key
    story.append(Paragraph("ANSWER KEY", heading_style))
    story.append(Paragraph("<b>Component Sorting:</b>", body_style))
    story.append(Paragraph("• EXTERNAL: Rainfall, Wind", bullet_style))
    story.append(Paragraph("• INTERNAL: Dry Vegetation, Fire Spread", bullet_style))
    story.append(Spacer(1, 8))

    story.append(Paragraph("<b>Relationships:</b>", body_style))
    story.append(Paragraph("• Rainfall → Dry Vegetation = <b>NEGATIVE</b> (more rain = less dry plants)", bullet_style))
    story.append(Paragraph("• Dry Vegetation → Fire Spread = <b>POSITIVE</b> (more fuel = more fire)", bullet_style))
    story.append(Paragraph("• Wind → Fire Spread = <b>POSITIVE</b> (more wind = fire spreads faster)", bullet_style))
    story.append(Spacer(1, 8))

    story.append(Paragraph("<b>Simulation Results:</b>", body_style))
    story.append(Paragraph("• Drought (Rainfall OFF): Dry Vegetation ↑, Fire Spread ↑", bullet_style))
    story.append(Paragraph("• Wind ON: Fire Spread ↑↑", bullet_style))
    story.append(Paragraph("• Worst conditions: No rain + High wind", bullet_style))
    story.append(Spacer(1, 15))

    # Page break
    story.append(PageBreak())

    # Facilitation Tips
    story.append(Paragraph("FACILITATION TIPS", heading_style))
    tips = [
        ("<b>Activity 1:</b>", "Let students explore. Ask: 'Which ones can we control?'"),
        ("<b>Activity 2:</b>", "Guide with: 'When this goes UP, does that go UP or DOWN?'"),
        ("<b>Activity 3:</b>", "Let students 'break' the model — turn things on/off. This is where insight happens!"),
        ("<b>Activity 4:</b>", "Don't give answers. Ask questions. Let curiosity drive research."),
    ]
    for label, tip in tips:
        story.append(Paragraph(f"{label} {tip}", bullet_style))
    story.append(Spacer(1, 15))

    # Discussion Prompts
    story.append(Paragraph("DISCUSSION PROMPTS", heading_style))
    prompts = [
        "\"Why does California burn at the same time every year?\"",
        "\"What would need to change to prevent these fires?\"",
        "\"How is this connected to the water cycle?\"",
        "\"If you were a firefighter, where would you position BEFORE fire season?\""
    ]
    for p in prompts:
        story.append(Paragraph(f"• {p}", bullet_style))
    story.append(Spacer(1, 15))

    # Common Misconceptions
    story.append(Paragraph("COMMON MISCONCEPTIONS", heading_style))
    misconceptions = [
        ("<b>\"Fires are caused by bad people\"</b>", "→ Most are human-caused, but conditions must be right. Without dry fuel, fires don't spread."),
        ("<b>\"More firefighters = no fires\"</b>", "→ Firefighters respond, but conditions (drought, wind) determine severity."),
        ("<b>\"Climate doesn't affect fires\"</b>", "→ Fire season has lengthened significantly due to climate patterns."),
    ]
    for label, response in misconceptions:
        story.append(Paragraph(f"{label} {response}", bullet_style))
    story.append(Spacer(1, 15))

    # Differentiation
    story.append(Paragraph("DIFFERENTIATION", heading_style))
    diff_data = [
        ['Support', 'Pre-label some relationships; use sentence starters'],
        ['Challenge', 'Add 3+ research components; calculate fire risk percentages'],
        ['ELL', 'Visual vocabulary cards; pair with English-proficient partner'],
    ]
    diff_table = Table(diff_data, colWidths=[1.2*inch, 5*inch])
    diff_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, -1), PDF_LIGHT_BLUE),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTNAME', (1, 0), (1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('GRID', (0, 0), (-1, -1), 0.5, PDF_MID_BLUE),
        ('PADDING', (0, 0), (-1, -1), 8),
    ]))
    story.append(diff_table)
    story.append(Spacer(1, 15))

    # STEM Challenge Notes
    story.append(Paragraph("STEM CHALLENGE: FIREBREAK ENGINEERS", heading_style))
    story.append(Paragraph("A simple, no-prep engineering challenge that extends the model.", body_style))
    story.append(Paragraph("• Students design firebreak placements to protect a community", bullet_style))
    story.append(Paragraph("• Connects directly to their model (reducing Dry Vegetation)", bullet_style))
    story.append(Paragraph("• Can be completed in 10-15 minutes", bullet_style))
    story.append(Paragraph("• Optional extension: Add 'Firebreaks' component to ModelIt", bullet_style))

    # Build PDF
    doc.build(story)
    print(f"[OK] Created: {output_path}")
    return output_path


# ============================================
# MAIN
# ============================================

if __name__ == "__main__":
    ensure_output_dir()
    print("\n" + "=" * 60)
    print("ModelIt! Branded Lesson Materials Generator")
    print("G05-L01: When Trees Become Matches")
    print("=" * 60 + "\n")

    ppt_path = create_powerpoint()
    activity_path = create_activity_pack()
    guide_path = create_teachers_guide()

    print("\n" + "=" * 60)
    print("All branded materials created!")
    print(f"Location: {OUTPUT_DIR}")
    print("=" * 60)
    print("\nFiles created:")
    print("  • G05-L01-Student-Presentation-BRANDED.pptx")
    print("  • G05-L01-Activity-Pack-BRANDED.pdf")
    print("  • G05-L01-Teachers-Guide-BRANDED.pdf")
    print("\nImage placeholders included - ready for NanoBanana generation!")
